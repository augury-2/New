"""
GARH KAUTHIG — Student Curtain-Raiser Deck  (vibrant, culture-rich)
===================================================================
A bright, celebratory event-introduction for the Swaragini Cultural Club's
"Garh Kauthig".  Built to excite students and to showcase the living culture of
Uttarakhand — folk dance, music, attire, food and craft.

Real event photos drop into the labelled PHOTO frames automatically (see
photos/README.md).  Run:

    python3 event.py           # -> GarhKauthig.pptx
"""
from __future__ import annotations
import math, os, random
from PIL import Image
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import design as D
import pptx_helpers as PX
from pptx_helpers import (text, picture, rect, hairline, dot, notes,
                          blank_slide, set_slide_bg)
from anim import Timeline, set_transition
import icons
import media

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
PHOTOS = os.path.join(HERE, "photos")
ICONDIR = os.path.join(ASSETS, "icons")

W, H = D.SLIDE_W_IN, D.SLIDE_H_IN
MX = 0.9
RR = MSO_SHAPE.ROUNDED_RECTANGLE
OV = MSO_SHAPE.OVAL


# ---------------------------------------------------------- vibrant palette --
class V:
    GREEN_DEEP = "0A3B2C"
    GREEN = "12583F"
    GREEN_BRIGHT = "1C7A54"
    CREAM = "FBF1DC"
    CREAM_DEEP = "F3E4C6"
    WHITE = "FFFDF7"
    INK = "12241C"
    SLATE = "3A4A42"

    MARIGOLD = "F5A81C"
    SAFFRON = "EE7B1E"
    RED = "E23A2E"
    MAGENTA = "D81E6A"
    TURQUOISE = "17A6A0"
    BLUE = "1E63B4"
    PURPLE = "7A3FB0"
    GOLD = "F2CE7B"
    GOLD_DEEP = "D9A43A"


ACCENTS = [V.MARIGOLD, V.RED, V.TURQUOISE, V.SAFFRON, V.MAGENTA, V.BLUE, V.PURPLE]
FD, FS, FE = D.F_DISPLAY, D.F_SERIF, D.F_EDITORIAL
FB, FU = D.F_BODY, D.F_UI

_prs = None


# ------------------------------------------------------------------ helpers --
def para(txt, **kw):
    kw["text"] = txt
    kw.setdefault("size", 16)
    kw.setdefault("color", V.WHITE)
    kw.setdefault("font", FB)
    return kw


def new_slide(bg=V.GREEN_DEEP, transition="morph", dur=D.D_TRANSITION):
    slide = blank_slide(_prs)
    set_slide_bg(slide, bg)
    set_transition(slide, transition, dur)
    return slide, Timeline(slide)


_ICON_CACHE = {}


def _icon_file(name, color):
    color = color.lstrip("#")
    key = (name, color)
    if key in _ICON_CACHE:
        return _ICON_CACHE[key]
    fn = f"ev_{name}_{color}.png"
    path = os.path.join(ICONDIR, fn)
    if not os.path.exists(path):
        try:
            icons.render(name, 240, color=color, lw=7.0).save(path)
        except Exception:
            _ICON_CACHE[key] = None
            return None
    _ICON_CACHE[key] = fn
    return fn


def eyebrow(slide, x, y, txt, color=V.GOLD, w=11.0):
    dot(slide, x + 0.06, y + 0.14, 0.13, fill=V.MARIGOLD)
    return [text(slide, x + 0.28, y, w, 0.34,
                 [para(txt, font=FU, size=12.5, color=color, caps=True,
                       track=300)], name="eyebrow")]


def icon_chip(slide, cx, cy, d, name, ring_color, glyph="12241C"):
    dot(slide, cx, cy, d, fill=ring_color)
    ic = _icon_file(name, glyph)
    if ic:
        s = d * 0.58
        picture(slide, os.path.join(ICONDIR, ic), cx - s / 2, cy - s / 2, s, s)


def footer(slide, label):
    text(slide, MX, H - 0.42, 7.0, 0.3,
         [para("गढ़ कौथिग  ·  GARH KAUTHIG", font=FU, size=9.5, color=V.GOLD,
               caps=True, track=180, cs=True)], name="fL")
    text(slide, W - MX - 5.0, H - 0.42, 5.0, 0.3,
         [para(label, font=FU, size=9.5, color=V.GOLD_DEEP, caps=True,
               track=160, align="r")], name="fR")


# ------------------------------------------------------- decorative motion ---
def _bleed_pic(slide, name, alpha):
    p = os.path.join(ASSETS, name)
    if not os.path.exists(p):
        return None
    pic = picture(slide, p, D.BLEED["x"], D.BLEED["y"], D.BLEED["w"],
                  D.BLEED["h"])
    PX.set_picture_transparency(pic, alpha)
    return pic


def living_bg(slide, tl, warm=True, rays=False, mandala=False, confetti=False,
              petals=True):
    """A layered, continuously-moving festival backdrop.

    Multiple independent layers drift/breathe at non-harmonic periods so the
    scene feels alive without ever distracting from the text."""
    # 1 — a big soft glow that slowly breathes
    glow = "glow_gold.png" if warm else "glow_cool.png"
    gp = os.path.join(ASSETS, glow)
    if os.path.exists(gp):
        g = picture(slide, gp, W * 0.5 - 5.2, -1.6, 10.4, 10.4)
        PX.set_picture_transparency(g, 0.34)
        tl.loop_zoom(g, 1.09, 11000)
    # 2 — a faint mandala watermark, breathing in a corner
    if mandala:
        mp = os.path.join(ASSETS, "orn_mandala_soft.png")
        if os.path.exists(mp):
            m = picture(slide, mp, W - 5.4, H - 5.2, 6.4, 6.4)
            PX.set_picture_transparency(m, 0.90)
            tl.loop_zoom(m, 1.05, 24000)
            tl.loop_drift(m, 0.006, 0.004, 21000)
    # 3 — light rays sweeping
    if rays:
        r = _bleed_pic(slide, "fx_rays_top.png", 0.55)
        if r is not None:
            tl.loop_drift(r, 0.02, 0.0, 19000)
    # 4 — two particle layers moving in different directions
    if petals:
        p1 = _bleed_pic(slide, "fx_petals.png", 0.20)
        if p1 is not None:
            tl.loop_drift(p1, -0.022, 0.030, 15000)
        p2 = _bleed_pic(slide, "fx_pollen.png", 0.34)
        if p2 is not None:
            tl.loop_drift(p2, 0.026, -0.020, 12000)
            tl.loop_zoom(p2, 1.05, 17000)
    # 5 — a few drifting colour specks (poster energy)
    if confetti:
        rng = random.Random(4)
        edge = [(0.06, 0.16), (0.16, 0.80), (0.9, 0.2), (0.85, 0.7),
                (0.5, 0.08), (0.28, 0.9), (0.72, 0.9), (0.95, 0.45)]
        for i, (fx, fy) in enumerate(edge):
            d = rng.uniform(0.09, 0.17)
            c = ACCENTS[i % len(ACCENTS)]
            sp = rect(slide, fx * W, fy * H, d, d, fill=c, alpha=0.55, shape=OV,
                      name="confetti")
            tl.loop_drift(sp, rng.uniform(-0.03, 0.03), rng.uniform(-0.04, 0.04),
                          rng.randint(9000, 16000))


def blobs(slide, specs):
    """Big soft colour blobs behind the content for a modern poster feel."""
    for (x, y, d, c, a) in specs:
        rect(slide, x, y, d, d, fill=c, alpha=a, shape=OV, name="blob")


def ornament_border(slide, top=True):
    p = os.path.join(ASSETS, "orn_aipan_gold.png")
    if os.path.exists(p):
        y = -0.15 if top else H - 0.72
        picture(slide, p, -0.1, y, W + 0.2, 0.88)


# ------------------------------------------------------------- photo frames --
def photo_frame(slide, x, y, w, h, label="YOUR PHOTO", accent=V.MARIGOLD,
                icon="mountain", tint=V.GREEN_BRIGHT, name=None):
    card = rect(slide, x, y, w, h, fill=tint, alpha=0.32, line=accent,
                line_w=2.6, shape=RR, radius=0.05, name=name or "photo")
    ic = _icon_file(icon, accent)
    if ic:
        s = min(w, h) * 0.30
        picture(slide, os.path.join(ICONDIR, ic), x + w / 2 - s / 2,
                y + h / 2 - s * 0.8, s, s)
    text(slide, x, y + h / 2 + min(w, h) * 0.14, w, 0.4,
         [para(label, font=FU, size=10.5, color=accent, caps=True, track=200,
               align="c")])
    return card


CROPDIR = os.path.join(PHOTOS, ".crop")


def _cover_crop(src, w, h):
    """Center-crop `src` to the w:h box aspect (no distortion) and cache it."""
    os.makedirs(CROPDIR, exist_ok=True)
    target = float(w) / float(h)
    base = os.path.splitext(os.path.basename(src))[0]
    out = os.path.join(CROPDIR, f"{base}_{target:.4f}.jpg")
    if os.path.exists(out) and os.path.getmtime(out) >= os.path.getmtime(src):
        return out
    im = Image.open(src).convert("RGB")
    iw, ih = im.size
    ar = iw / ih
    if ar > target:                      # too wide -> trim the sides
        nw = int(round(ih * target))
        x0 = (iw - nw) // 2
        im = im.crop((x0, 0, x0 + nw, ih))
    elif ar < target:                    # too tall -> trim top/bottom
        nh = int(round(iw / target))
        y0 = (ih - nh) // 2
        im = im.crop((0, y0, iw, y0 + nh))
    im.save(out, "JPEG", quality=90)
    return out


def real_or_frame(slide, x, y, w, h, filename, crop=True, border=True,
                  radius=0.04, **kw):
    for ext in ("", ".jpg", ".jpeg", ".png"):
        p = os.path.join(PHOTOS, filename if ext == "" else
                         os.path.splitext(filename)[0] + ext)
        if os.path.exists(p):
            src = _cover_crop(p, w, h) if crop else p
            pic = picture(slide, src, x, y, w, h)
            if border:
                acc = kw.get("accent", V.GOLD)
                rect(slide, x, y, w, h, fill=None, line=acc, line_w=2.2,
                     shape=RR, radius=radius, name="photo_edge")
            return pic
    return photo_frame(slide, x, y, w, h, **kw)


# --------------------------------------------------------------- card grids --
def culture_card(slide, x, y, w, h, icon, name, desc, accent, on_dark=True):
    fill = V.WHITE if on_dark else V.GREEN_DEEP
    name_c = V.GREEN_DEEP if on_dark else V.WHITE
    desc_c = V.SLATE if on_dark else V.GOLD
    c = rect(slide, x, y, w, h, fill=fill, alpha=1.0, line=accent, line_w=1.6,
             shape=RR, radius=0.06, name="cc")
    rect(slide, x, y, w, 0.12, fill=accent, alpha=1.0, shape=RR, radius=0.5,
         name="ccbar")
    chip_d = min(0.8, h * 0.34)
    icon_chip(slide, x + w / 2, y + 0.28 + chip_d / 2, chip_d, icon, accent)
    name_y = y + 0.34 + chip_d
    text(slide, x + 0.08, name_y, w - 0.16, 0.56,
         [para(name, font=FE, size=15, color=name_c, align="c", line=1.0)],
         name="ccname")
    desc_y = name_y + 0.56
    text(slide, x + 0.16, desc_y, w - 0.32, y + h - desc_y - 0.1,
         [para(desc, font=FB, size=10.5, color=desc_c, align="c", line=1.22)],
         name="ccdesc")
    return c


def culture_slide(bg, on_dark, eb_txt, title_specs, items, cols, foot, note,
                  warm=True, rays=False):
    slide, tl = new_slide(bg=bg)
    living_bg(slide, tl, warm=warm, mandala=True, rays=rays, petals=True)
    eb = eyebrow(slide, MX, 0.66, eb_txt, color=(V.GOLD if on_dark else V.SAFFRON))
    ttl = text(slide, MX, 1.06, W - 2 * MX, 0.9, title_specs, name="title")

    n = len(items)
    rows = math.ceil(n / cols)
    gx, gy = 0.28, 0.28
    top = 2.18
    avail_h = (H - 0.52) - top
    ch = min(2.75, (avail_h - (rows - 1) * gy) / rows)
    cw = (W - 2 * MX - (cols - 1) * gx) / cols
    shapes = []
    for i, (ic, nm, ds, ac) in enumerate(items):
        r, cidx = divmod(i, cols)
        # centre a short last row
        in_row = min(cols, n - r * cols)
        row_w = in_row * cw + (in_row - 1) * gx
        x0 = (W - row_w) / 2
        x = x0 + cidx * (cw + gx)
        yy = top + r * (ch + gy)
        shapes.append(culture_card(slide, x, yy, cw, ch, ic, nm, ds, ac,
                                   on_dark=on_dark))
    for s in eb:
        tl.fade_in(s, delay=200)
    tl.rise_in(ttl, delay=450, dur=1200)
    tl.stagger(shapes, start=900, gap=150, mode="rise")
    footer(slide, foot)
    tl.apply()
    notes(slide, note)


# ============================================================= SLIDES ========
def s_cover():
    slide, tl = new_slide(bg=V.GREEN_DEEP, transition="none")
    blobs(slide, [(-1.5, -1.6, 4.5, V.MARIGOLD, 0.14),
                  (W - 3.2, -2.0, 5.2, V.MAGENTA, 0.12),
                  (W - 2.2, H - 2.6, 4.0, V.TURQUOISE, 0.12),
                  (-1.8, H - 3.0, 4.4, V.SAFFRON, 0.12)])
    living_bg(slide, tl, warm=True, rays=True, confetti=True, petals=True)
    ornament_border(slide, top=True)

    eb = eyebrow(slide, MX, 0.98, "Swaragini Cultural Club  ·  Presents",
                 color=V.GOLD)
    hi = text(slide, MX, 1.42, W - 2 * MX, 2.2,
              [para("गढ़ कौथिग", font=FD, size=120, color=V.MARIGOLD, bold=True,
                    cs=True, shadow=dict(blur=28, dist=6, alpha=60))],
              name="hindi")
    sub = text(slide, MX, 3.86, W - 2 * MX, 1.0,
               [para("A Celebration of Uttarakhand\u2019s Culture", font=FE,
                     size=40, color=V.WHITE,
                     shadow=dict(blur=18, dist=3, alpha=50))], name="sub")
    r = hairline(slide, MX, 4.95, 3.2, color=V.MARIGOLD, weight=2.6)
    tag = text(slide, MX, 5.12, W - 2 * MX, 0.5,
               [para("Dance  ·  Music  ·  Food  ·  Craft  ·  Community",
                     font=FU, size=15, color=V.GOLD, caps=True, track=240)],
               name="tag")

    files = ["cover1", "cover2", "cover3", "cover4", "cover5"]
    ic5 = ["dancer", "temple", "dhol", "mountain", "hands"]
    fw = (W - 2 * MX - 4 * 0.18) / 5
    frames = []
    for i, fn in enumerate(files):
        fx = MX + i * (fw + 0.18)
        frames.append(real_or_frame(slide, fx, 5.92, fw, 1.06, fn,
                      label="PHOTO", accent=ACCENTS[i % len(ACCENTS)],
                      icon=ic5[i], tint=V.GREEN_BRIGHT))

    aud = media.add_audio(slide, os.path.join(ASSETS, "score.m4a"),
                          x=W - 0.55, y=0.26, size=0.40)
    for s in eb:
        tl.fade_in(s, delay=300)
    tl.rise_in(hi, delay=600, dur=1500, rise=0.05)
    tl.rise_in(sub, delay=1300, dur=1400)
    tl.wipe_in(r, delay=1900, dur=800)
    tl.fade_in(tag, delay=2150, dur=1200)
    tl.stagger(frames, start=2400, gap=170, mode="rise")
    tl.play_media(aud, loop=True, vol=60000, kind="audio")
    footer(slide, "Curtain Raiser")
    tl.apply()
    notes(slide, "Open with energy. Let the dhol-damau track play. This is a "
                 "surprise curtain-raiser \u2014 build excitement from slide 1.")


def s_represents():
    slide, tl = new_slide(bg=V.GREEN_DEEP)
    living_bg(slide, tl, warm=True, mandala=True, petals=True)
    eb = eyebrow(slide, MX, 0.82, "What It Means")
    ttl = text(slide, MX, 1.24, 6.6, 1.6,
               [para("Garh Kauthig", font=FE, size=52, color=V.MARIGOLD),
                para("means celebration", font=FE, size=40, color=V.WHITE)],
               name="title")
    r = hairline(slide, MX, 3.2, 2.4, color=V.MARIGOLD, weight=2.2)
    body = text(slide, MX, 3.45, 6.1, 2.8,
                [para("\u201cGarh Kauthik\u201d is a traditional Uttarakhandi "
                      "word for a fair, festival or community celebration.",
                      font=FS, size=23, color=V.WHITE, italic=True, line=1.3),
                 para("It brings people together to celebrate the region\u2019s "
                      "culture, music, dance, food, clothing and traditions.",
                      font=FB, size=16, color=V.GOLD, line=1.4, before=14)],
                name="body")
    fr = real_or_frame(slide, 7.45, 1.35, 4.95, 4.9, "represents",
                       label="FESTIVAL PHOTO", accent=V.MARIGOLD, icon="dhol",
                       tint=V.GREEN_BRIGHT)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.rise_in(ttl, delay=500, dur=1300)
    tl.wipe_in(r, delay=1150, dur=700)
    tl.rise_in(body, delay=1350, dur=1300)
    tl.zoom_in(fr, delay=900, dur=1500)
    footer(slide, "The Meaning")
    tl.apply()
    notes(slide, "Kauthig = a fair / celebration that gathers the whole "
                 "community. Keep it simple and warm.")


def s_where():
    slide, tl = new_slide(bg=V.CREAM)
    rect(slide, -0.1, 0, W + 0.2, 2.3, fill=V.GREEN_DEEP, alpha=1.0)
    living_bg(slide, tl, warm=True, mandala=True, petals=False, confetti=True)
    eb = eyebrow(slide, MX, 0.62, "Across The Hills", color=V.GOLD)
    ttl = text(slide, MX, 1.0, W - 2 * MX, 0.7,
               [para("Where is Garh Kauthig celebrated?", font=FE, size=33,
                     color=V.WHITE)], name="title")
    lead = text(slide, MX, 1.62, W - 2 * MX, 0.5,
                [para("In villages, towns and cities across Uttarakhand \u2014 "
                      "especially Garhwal. It comes alive during:", font=FB,
                      size=14, color=V.GOLD, align="c")], name="lead")
    items = [
        ("seedling", "HARVEST FESTIVALS",
         "Gratitude for bountiful crops and prosperity."),
        ("temple", "SACRED DAYS",
         "Rituals and spiritual days across the hamlets."),
        ("cart", "COMMUNITY MELAS",
         "Seasonal harvest gatherings and markets."),
        ("flag", "TEMPLE FAIRS",
         "Deity processions and thanksgiving melas."),
        ("youth", "MODERN STAGES",
         "Youth platforms celebrating folk arts today."),
    ]
    cw = (W - 2 * MX - 4 * 0.3) / 5
    ch = 3.05
    # centre the card row within the cream band below the header
    top = 2.3 + ((H - 0.5) - 2.3 - ch) / 2
    cards = []
    for i, (ic, hd, bd) in enumerate(items):
        x = MX + i * (cw + 0.3)
        acc = ACCENTS[i % len(ACCENTS)]
        card = rect(slide, x, top, cw, ch, fill=V.WHITE, alpha=1.0, line=acc,
                    line_w=1.6, shape=RR, radius=0.06, name="wcard")
        rect(slide, x, top, cw, 0.12, fill=acc, alpha=1.0, shape=RR,
             radius=0.5, name="wbar")
        icon_chip(slide, x + cw / 2, top + 0.66, 0.78, ic, acc)
        hh = text(slide, x + 0.05, top + 1.22, cw - 0.1, 0.6,
                  [para(hd, font=FU, size=11.5, color=V.GREEN_DEEP, bold=True,
                        caps=True, track=20, align="c", line=1.05)], name="wh")
        bb = text(slide, x + 0.14, top + 1.86, cw - 0.28, ch - 2.0,
                  [para(bd, font=FB, size=10.5, color=V.SLATE, align="c",
                        line=1.26)], name="wb")
        cards += [card, hh, bb]
    for s in eb:
        tl.fade_in(s, delay=200)
    tl.fade_in(ttl, delay=450)
    tl.fade_in(lead, delay=700)
    tl.stagger(cards, start=1000, gap=130, mode="rise")
    footer(slide, "Where")
    tl.apply()
    notes(slide, "Five occasions. Emphasise it is woven into farming, faith and "
                 "now student life.")


def s_happens():
    slide, tl = new_slide(bg=V.GREEN_DEEP)
    living_bg(slide, tl, warm=True, mandala=True, petals=True)
    eb = eyebrow(slide, MX, 0.82, "The Experience")
    ttl = text(slide, MX, 1.22, 6.4, 1.6,
               [para("What happens at a", font=FE, size=38, color=V.WHITE),
                para("Garh Kauthig", font=FE, size=44, color=V.MARIGOLD)],
               name="title")
    acts = [("dancer", "Garhwali & Kumaoni folk dance and music"),
            ("flag", "Cultural competitions and regional games"),
            ("sweet", "Local, authentic food stalls and handicrafts"),
            ("jewel", "Traditional clothing and mountain jewellery")]
    rows = []
    for i, (ic, tx) in enumerate(acts):
        y = 2.95 + i * 0.82
        acc = ACCENTS[i % len(ACCENTS)]
        icon_chip(slide, MX + 0.34, y + 0.24, 0.64, ic, acc)
        rows.append(text(slide, MX + 0.9, y, 5.5, 0.7,
                    [para(tx, font=FB, size=16, color=V.WHITE, line=1.12)],
                    name="act"))
    panel = rect(slide, 7.2, 1.45, 5.2, 4.65, fill=V.MARIGOLD, alpha=1.0,
                 shape=RR, radius=0.05, name="panel")
    ph = text(slide, 7.55, 1.8, 4.55, 0.6,
              [para("POPULAR FOLK PERFORMANCES", font=FU, size=13,
                    color=V.GREEN_DEEP, bold=True, caps=True, track=120)],
              name="ph")
    pv = text(slide, 7.55, 2.5, 4.55, 1.3,
              [para("Jhora · Chanchari · Chholiya", font=FE, size=20,
                    color=V.INK, line=1.22),
               para("Pandav Nritya · Langvir Nritya", font=FE, size=20,
                    color=V.INK, line=1.22, before=6)], name="pv")
    pb = text(slide, 7.55, 3.95, 4.55, 1.9,
              [para("Mesmerising performances that carry the martial spirit and "
                    "living culture of the mountains \u2014 danced in a circle, "
                    "together.", font=FB, size=13, color=V.GREEN_DEEP,
                    line=1.34)], name="pb")
    for s in eb:
        tl.fade_in(s, delay=200)
    tl.rise_in(ttl, delay=450, dur=1300)
    tl.stagger(rows, start=1000, gap=200, mode="rise")
    tl.zoom_in(panel, delay=900, dur=1200)
    tl.fade_in(ph, delay=1400)
    tl.rise_in(pv, delay=1650, dur=1200)
    tl.fade_in(pb, delay=2050)
    footer(slide, "What Happens")
    tl.apply()
    notes(slide, "Dance, music, competitions, food, crafts, costume. Name the "
                 "folk forms \u2014 students may get to learn some.")


# ---- culture section --------------------------------------------------------
def s_showcase():
    """Full-bleed, looping real-photo dance montage with the dhol-damau score."""
    slide, tl = new_slide(bg=V.INK, transition="fade")
    vid = os.path.join(ASSETS, "montage.mp4")
    poster = os.path.join(ASSETS, "montage_poster.jpg")
    movie = None
    if os.path.exists(vid):
        movie = media.add_video(slide, vid, poster, 0, 0, W, H)
    else:
        # graceful fallback if the video has not been rendered yet
        real_or_frame(slide, 0, 0, W, H, "cover2", crop=True, border=False,
                      accent=V.MARIGOLD, icon="dancer")

    # legibility scrim across the lower third
    rect(slide, -0.1, H - 3.1, W + 0.2, 0.9, fill=V.INK, alpha=0.30,
         name="scrim_soft")
    rect(slide, -0.1, H - 2.25, W + 0.2, 2.4, fill=V.INK, alpha=0.55,
         name="scrim")

    eb = eyebrow(slide, MX, H - 1.98, "Our Students, Live On Stage",
                 color=V.GOLD)
    ttl = text(slide, MX, H - 1.6, W - 2 * MX, 1.0,
               [para("Feel the rhythm of the hills", font=FE, size=44,
                     color=V.WHITE, shadow=dict(blur=20, dist=3, alpha=55))],
               name="title")
    sub = text(slide, MX, H - 0.86, W - 2 * MX, 0.4,
               [para("Real moments from Garh Kauthig \u2014 danced to the "
                     "dhol-damau beat.", font=FB, size=14.5, color=V.GOLD)],
               name="sub")
    for s in eb:
        tl.fade_in(s, delay=400)
    tl.rise_in(ttl, delay=650, dur=1300)
    tl.fade_in(sub, delay=1200, dur=1200)
    if movie is not None:
        tl.play_media(movie, loop=True, vol=85000, kind="video")
    footer(slide, "In Motion")
    tl.apply()
    notes(slide, "Let it play. This is real footage of our own students "
                 "performing at Garh Kauthig, set to the dhol-damau score \u2014 "
                 "the surprise energy hit of the curtain-raiser.")


def s_dances():
    culture_slide(
        bg=V.GREEN_DEEP, on_dark=True, eb_txt="Living Culture  ·  The Dance",
        title_specs=[para("Folk dances that tell our stories", font=FE, size=34,
                          color=V.WHITE)],
        items=[
            ("jhora", "Jhora", "Everyone in one circle \u2014 arms linked, "
             "swaying to the hurka drum, across every community.", V.MARIGOLD),
            ("sword", "Chholiya", "The sword dance: shields and blades flash to "
             "guard a wedding procession from evil.", V.RED),
            ("hands", "Chanchari", "A gentle group circle from the Danpur "
             "valley, danced slow and together.", V.TURQUOISE),
            ("trident", "Pandav Nritya", "The Mahabharata, danced \u2014 the "
             "Pandavas\u2019 tale through trance, song and drum.", V.PURPLE),
            ("dancer", "Langvir Nritya", "The daring one: a lone acrobat climbs "
             "a tall pole and balances at the top.", V.BLUE),
        ], cols=5, foot="Folk Dance",
        note="Five signature dances. Jhora = togetherness; Chholiya = martial "
             "sword dance; Pandav Nritya = epic re-enactment; Langvir = "
             "acrobatic. Students can try Jhora and Chholiya.")


def s_music():
    culture_slide(
        bg=V.CREAM, on_dark=False, eb_txt="Living Culture  ·  The Sound",
        title_specs=[para("The instruments that set the beat", font=FE, size=34,
                          color=V.GREEN_DEEP)],
        items=[
            ("dhol", "Dhol", "The big double-headed drum \u2014 the heartbeat "
             "of every fair.", V.RED),
            ("damau", "Damau", "The dhol\u2019s copper partner, keeping the "
             "fast pulse.", V.MARIGOLD),
            ("ransingha", "Ransingha", "A curved brass horn whose call carries "
             "for miles.", V.SAFFRON),
            ("bhankora", "Bhankora", "A long straight temple trumpet of "
             "Garhwal.", V.TURQUOISE),
            ("hurka", "Hurka", "The hourglass drum that leads the Jhora "
             "circle.", V.BLUE),
            ("masakbeen", "Masakbeen", "The Himalayan bagpipe \u2014 a reed for "
             "celebration.", V.MAGENTA),
        ], cols=3, foot="Folk Music", warm=True,
        note="Dhol-damau are the soul of the hills. Point out that the "
             "background track you can hear is a dhol-damau rhythm.")


def s_attire():
    culture_slide(
        bg=V.GREEN_DEEP, on_dark=True, eb_txt="Living Culture  ·  What We Wear",
        title_specs=[para("Colour, cloth and mountain gold", font=FE, size=34,
                          color=V.WHITE)],
        items=[
            ("pichhora", "Rangwali Pichhora", "The hand-painted yellow odhni of "
             "married women, made while singing blessings.", V.MARIGOLD),
            ("loom", "Ghagra · Choli · Orni", "A flowing skirt, fitted blouse "
             "and veil \u2014 the everyday grace of the hills.", V.MAGENTA),
            ("jewel", "Nath", "The grand golden nose-ring, the pride of a "
             "Kumaoni bride.", V.RED),
            ("coin", "Hansuli & Guloband", "A silver neck-torque and a "
             "red-velvet choker \u2014 mountain heirlooms.", V.TURQUOISE),
        ], cols=4, foot="Traditional Attire",
        note="Pichhora is the signature yellow bridal veil. Jewellery like nath "
             "and hansuli are heirlooms passed down generations.")


def s_food():
    culture_slide(
        bg=V.CREAM, on_dark=False, eb_txt="Living Culture  ·  Flavours",
        title_specs=[para("Flavours of the hills", font=FE, size=34,
                          color=V.GREEN_DEEP)],
        items=[
            ("thali", "Aloo ke Gutke", "Spiced mountain potatoes with wild "
             "jakhiya seed.", V.RED),
            ("leaf", "Kafuli", "A creamy green curry of spinach and "
             "fenugreek.", V.TURQUOISE),
            ("millet", "Mandua Roti", "Soft finger-millet bread \u2014 the hill "
             "staple.", V.SAFFRON),
            ("sweet", "Bal Mithai", "Almora\u2019s chocolate-brown fudge with "
             "sugar pearls.", V.MAGENTA),
            ("pot", "Bhatt · Chainsoo", "Hearty roasted-lentil dals, full of "
             "protein.", V.BLUE),
            ("chilli", "Bhang ki Chutney", "A tangy roasted hemp-seed chutney "
             "of the hills.", V.PURPLE),
        ], cols=3, foot="Cuisine", warm=True,
        note="Millets (mandua, jhangora), potatoes, lentils and local greens. "
             "Bal Mithai and Singori are the famous sweets.")


def s_crafts():
    culture_slide(
        bg=V.GREEN_DEEP, on_dark=True, eb_txt="Living Culture  ·  Made By Hand",
        title_specs=[para("Handmade heritage of Uttarakhand", font=FE, size=34,
                          color=V.WHITE)],
        items=[
            ("aipan", "Aipan", "White rice-paste patterns on red earth, drawn "
             "by women for every blessing.", V.RED),
            ("ringal", "Ringaal", "Hill bamboo woven into baskets, mats and "
             "rings.", V.MARIGOLD),
            ("copper", "Tamta Copper", "Almora\u2019s coppersmiths beat gleaming "
             "pots and temple bells by hand.", V.SAFFRON),
            ("chisel", "Woodcarving", "Ornate carved doorways \u2014 kholi \u2014 "
             "frame the old hill homes.", V.TURQUOISE),
            ("loom", "Wool Weaving", "Warm carpets, blankets and shawls spun "
             "from mountain wool.", V.MAGENTA),
        ], cols=5, foot="Handicrafts",
        note="Aipan is the ritual folk art of Kumaon. Ringaal, tamta copper, "
             "woodcarving and wool are living crafts, several now GI-tagged.")


def s_offer():
    slide, tl = new_slide(bg=V.CREAM)
    living_bg(slide, tl, warm=False, mandala=True, petals=False, confetti=True)
    eb = eyebrow(slide, MX, 0.8, "For Every Participant", color=V.SAFFRON)
    ttl = text(slide, MX, 1.2, W - 2 * MX, 0.9,
               [para("What will this event offer you?", font=FE, size=40,
                     color=V.GREEN_DEEP)], name="title")
    lead = text(slide, MX, 2.12, 10, 0.5,
                [para("A chance to grow \u2014 culturally, creatively and "
                      "personally:", font=FB, size=15, color=V.SAFFRON)],
                name="lead")
    offers = [("book", "Learn about Uttarakhand\u2019s rich culture"),
              ("hands", "Develop creative and performance skills"),
              ("dhol", "Experience traditional music and dance"),
              ("sunrise", "An exclusive stage to showcase your talent"),
              ("youth", "Build confidence, teamwork and cultural roots")]
    cw = (W - 2 * MX - 0.4) / 2
    cards = []
    for i, (ic, tx) in enumerate(offers):
        col, row = i % 2, i // 2
        x = MX + col * (cw + 0.4)
        y = 2.85 + row * 1.24
        if i == 4:
            x = MX + (W - 2 * MX) / 2 - cw / 2
        acc = ACCENTS[i % len(ACCENTS)]
        c = rect(slide, x, y, cw, 1.04, fill=V.WHITE, alpha=1.0, line=acc,
                 line_w=1.6, shape=RR, radius=0.09, name="ocard")
        icon_chip(slide, x + 0.6, y + 0.52, 0.66, ic, acc)
        tt = text(slide, x + 1.05, y, cw - 1.25, 1.04,
                  [para(tx, font=FB, size=15, color=V.GREEN_DEEP, line=1.14)],
                  name="ot", anchor=MSO_ANCHOR.MIDDLE)
        cards += [c, tt]
    for s in eb:
        tl.fade_in(s, delay=200)
    tl.rise_in(ttl, delay=450, dur=1200)
    tl.fade_in(lead, delay=800)
    tl.stagger(cards, start=1100, gap=130, mode="rise")
    footer(slide, "What You Gain")
    tl.apply()
    notes(slide, "Sell it to students: culture, skills, stage time, confidence, "
                 "friendships.")


def s_guide():
    slide, tl = new_slide(bg=V.GREEN_DEEP)
    living_bg(slide, tl, warm=True, mandala=True, petals=True)
    eb = eyebrow(slide, MX, 0.82, "Your Mentors")
    ttl = text(slide, MX, 1.24, 6.2, 1.0,
               [para("Who will guide you?", font=FE, size=42, color=V.WHITE)],
               name="title")
    body = text(slide, MX, 2.5, 5.6, 2.8,
                [para("Our teams", font=FB, size=17, color=V.GOLD),
                 para("DevBhoomi & Uttarayani", font=FE, size=32,
                      color=V.MARIGOLD, line=1.12, before=4),
                 para("will guide participants and teach them about "
                      "Uttarakhand\u2019s traditional culture, folk "
                      "performances, and authentic artistic presentation.",
                      font=FB, size=16, color=V.WHITE, line=1.4, before=12)],
                name="body")
    f1 = real_or_frame(slide, 6.85, 1.45, 2.6, 2.95, "team1",
                       label="TEAM DEVBHOOMI", accent=V.RED, icon="dancer",
                       tint=V.GREEN_BRIGHT)
    f2 = real_or_frame(slide, 9.72, 1.45, 2.68, 2.95, "team2",
                       label="TEAM UTTARAYANI", accent=V.TURQUOISE,
                       icon="hands", tint=V.GREEN_BRIGHT)
    f3 = real_or_frame(slide, 6.85, 4.6, 5.55, 2.1, "team3",
                       label="PERFORMANCE PHOTO", accent=V.MARIGOLD,
                       icon="flag", tint=V.GREEN_BRIGHT)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.rise_in(ttl, delay=500, dur=1300)
    tl.rise_in(body, delay=900, dur=1300)
    tl.stagger([f1, f2, f3], start=1000, gap=220, mode="zoom")
    footer(slide, "Your Guides")
    tl.apply()
    notes(slide, "Introduce mentor teams DevBhoomi and Uttarayani \u2014 the "
                 "seniors who will train participants.")


def s_win():
    slide, tl = new_slide(bg=V.GREEN_DEEP)
    blobs(slide, [(-1.6, -1.8, 4.6, V.MARIGOLD, 0.13),
                  (W - 3.0, H - 3.0, 5.0, V.RED, 0.12)])
    living_bg(slide, tl, warm=True, rays=True, confetti=True, petals=True)
    ornament_border(slide, top=False)
    eb = eyebrow(slide, MX, 0.98, "Your Moment On Stage", color=V.GOLD)
    ttl = text(slide, MX, 1.45, W - 2 * MX, 1.6,
               [para("Showcase your talent \u2014 and win!", font=FE, size=48,
                     color=V.WHITE, shadow=dict(blur=18, dist=3, alpha=45))],
               name="title")
    body = text(slide, MX, 2.95, W - 2 * MX, 0.9,
                [para("Present the true beauty of Uttarakhand\u2019s culture "
                      "through your own creativity and performance.", font=FB,
                      size=18, color=V.GOLD, align="c", line=1.35)],
                name="body")
    pill = rect(slide, W / 2 - 2.7, 4.1, 5.4, 2.15, fill=V.MARIGOLD, alpha=1.0,
                shape=RR, radius=0.11, name="pill")
    plabel = text(slide, W / 2 - 2.7, 4.38, 5.4, 0.5,
                  [para("CASH PRIZE", font=FU, size=17, color=V.GREEN_DEEP,
                        bold=True, caps=True, track=280, align="c")],
                  name="plabel")
    prize = text(slide, W / 2 - 2.7, 4.86, 5.4, 1.2,
                 [para("\u20B9 11,000", font=FD, size=74, color=V.RED,
                       bold=True, align="c")], name="prize")
    for s in eb:
        tl.fade_in(s, delay=300)
    tl.rise_in(ttl, delay=650, dur=1400)
    tl.fade_in(body, delay=1300)
    tl.zoom_in(pill, delay=1600, dur=1100)
    tl.fade_in(plabel, delay=2000)
    tl.rise_in(prize, delay=2200, dur=1200)
    footer(slide, "Compete")
    tl.apply()
    notes(slide, "The hook: participate and win. Announce the \u20B911,000 cash "
                 "prize with energy \u2014 this is the call to action.")


def s_closing():
    slide, tl = new_slide(bg=V.CREAM, transition="smooth",
                          dur=D.D_TRANSITION_SLOW)
    living_bg(slide, tl, warm=True, mandala=True, confetti=True, petals=True)
    ornament_border(slide, top=True)
    ornament_border(slide, top=False)
    ln = text(slide, 0.6, 2.5, W - 1.2, 2.6,
              [para("Celebrate our traditions,", font=FE, size=42,
                    color=V.GREEN_DEEP, align="c", line=1.22),
               para("showcase your talent, and keep the", font=FE, size=42,
                    color=V.GREEN_DEEP, align="c", line=1.22),
               para("culture of Uttarakhand alive.", font=FE, size=42,
                    color=V.RED, align="c", line=1.22, before=6)], name="line")
    tl.rise_in(ln, delay=500, dur=1800, rise=0.03)
    footer(slide, "The Invitation")
    tl.apply()
    notes(slide, "Land the line, then invite them in. Pause before Thank You.")


def s_thanks():
    slide, tl = new_slide(bg=V.GREEN_DEEP, transition="smooth",
                          dur=D.D_TRANSITION_SLOW)
    blobs(slide, [(-1.6, -1.8, 4.4, V.MAGENTA, 0.13),
                  (W - 3.0, H - 3.0, 5.0, V.TURQUOISE, 0.12)])
    living_bg(slide, tl, warm=True, confetti=True, petals=True)
    ornament_border(slide, top=True)
    ornament_border(slide, top=False)
    ty = text(slide, MX, 2.35, W - 2 * MX, 1.6,
              [para("THANK YOU!", font=FD, size=90, color=V.MARIGOLD, bold=True,
                    align="c", shadow=dict(blur=24, dist=4, alpha=55))],
              name="ty")
    sub = text(slide, MX, 4.5, W - 2 * MX, 0.7,
               [para("Presented by the Swaragini Cultural Club", font=FU,
                     size=16, color=V.GOLD, caps=True, track=200, align="c")],
               name="sub")
    tl.rise_in(ty, delay=500, dur=1500, rise=0.04)
    tl.fade_in(sub, delay=1300, dur=1300)
    tl.apply()
    notes(slide, "Close with warmth. Invite questions and sign-ups.")


# ================================================================= build =====
def build():
    global _prs
    _prs = Presentation()
    _prs.slide_width = D.SLIDE_W_EMU
    _prs.slide_height = D.SLIDE_H_EMU

    s_cover()
    s_represents()
    s_where()
    s_happens()
    s_showcase()
    s_dances()
    s_music()
    s_attire()
    s_food()
    s_crafts()
    s_offer()
    s_guide()
    s_win()
    s_closing()
    s_thanks()

    out = os.path.join(HERE, "GarhKauthig.pptx")
    _prs.save(out)
    print(f"saved {out}  ({len(_prs.slides._sldIdLst)} slides, "
          f"{os.path.getsize(out)/1e6:.1f} MB)")


if __name__ == "__main__":
    os.makedirs(PHOTOS, exist_ok=True)
    build()
