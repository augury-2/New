"""
KAUTHIK — Media embedding
=========================
Thin helpers to embed the (original, procedurally generated) dance video and
ambient audio into the deck.  python-pptx ships `add_movie` for video but has
no audio equivalent, so the audio picture is built by hand following the same
OOXML pattern PowerPoint uses (`a:audioFile` + `p14:media` + a poster image).
"""
from __future__ import annotations
import io, os
from PIL import Image, ImageDraw

from pptx.util import Emu
from pptx.oxml.ns import nsdecls, qn
from pptx.oxml import parse_xml
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.media import Video

import design as D
from design import rgb, rgba

HERE = os.path.dirname(os.path.abspath(__file__))
IN = D.EMU_PER_IN


def _emu(v):
    return Emu(int(round(v * IN)))


# ------------------------------------------------------------------- video ---
def add_video(slide, path, poster, x, y, w, h, mime="video/mp4"):
    """Embed a looping background video; returns the movie shape."""
    shp = slide.shapes.add_movie(path, _emu(x), _emu(y), _emu(w), _emu(h),
                                 poster_frame_image=poster, mime_type=mime)
    return shp


# ------------------------------------------------------------------- audio ---
def _speaker_poster() -> io.BytesIO:
    """A small, tasteful gold speaker glyph used as the audio placeholder."""
    S = 220
    img = Image.new("RGBA", (S, S), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.ellipse([6, 6, S - 6, S - 6], fill=rgba(D.C.INK, 235),
              outline=rgba(D.C.GOLD, 235), width=6)
    g = rgb(D.C.GOLD_LIGHT)
    # speaker body
    d.polygon([(70, 92), (100, 92), (128, 66), (128, 154), (100, 128), (70, 128)],
              fill=g)
    # sound waves
    for r, wdt in ((150, 8), (176, 7)):
        d.arc([128 - (r - 128), 110 - (r - 60), 128 + (r - 128), 110 + (r - 60)],
              -55, 55, fill=g, width=wdt)
    buf = io.BytesIO()
    img.save(buf, "PNG")
    buf.seek(0)
    return buf


def add_audio(slide, path, x, y, size=0.42, mime="audio/mp4"):
    """
    Embed an audio clip as a small speaker icon.  Returns the `p:pic` shape
    element (has `.shape_id`).  Autoplay/loop is applied separately through
    the Timeline's `play_media(..., kind='audio')`.
    """
    part = slide.part
    video = Video.from_path_or_file_like(path, mime)
    media_part = part._package.get_or_add_media_part(video)
    media_rId = part.relate_to(media_part, RT.MEDIA)
    audio_rId = part.relate_to(media_part, RT.AUDIO)
    _, poster_rId = part.get_or_add_image_part(_speaker_poster())

    sid = slide.shapes._next_shape_id
    xml = (
        '<p:pic %s>'
        '<p:nvPicPr>'
        '<p:cNvPr id="%d" name="Ambient Score">'
        '<a:hlinkClick r:id="" action="ppaction://media"/>'
        '</p:cNvPr>'
        '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
        '<p:nvPr>'
        '<a:audioFile r:link="%s"/>'
        '<p:extLst>'
        '<p:ext uri="{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}">'
        '<p14:media xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
        'r:embed="%s"/>'
        '</p:ext></p:extLst>'
        '</p:nvPr>'
        '</p:nvPicPr>'
        '<p:blipFill><a:blip r:embed="%s"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
        '<p:spPr>'
        '<a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
        '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        '</p:spPr>'
        '</p:pic>'
        % (nsdecls("a", "p", "r"), sid, audio_rId, media_rId, poster_rId,
           int(round(x * IN)), int(round(y * IN)),
           int(round(size * IN)), int(round(size * IN)))
    )
    pic = parse_xml(xml)
    slide.shapes._spTree.append(pic)
    return slide.shapes._shape_factory(pic)
