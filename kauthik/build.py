"""
KAUTHIK — Deck Builder
======================
Assembles the full presentation: ~30 cinematic 16:9 slides, each with a living
parallax background, choreographed foreground entrances and a Morph transition
into the next. Run as a script to (re)generate everything and write the
`.pptx`.

    python3 build.py               # incremental (reuses cached art)
    python3 build.py --force-art   # regenerate all artwork first
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

import design as D
from design import C
import content as K
import icons
import mapgen
import media
import scenes
from anim import Timeline, set_transition
from pptx_helpers import (IN, blank_slide, bring_to_front, dot, hairline,
                          picture, rect, ring, notes, send_to_back,
                          set_slide_bg, text, vline)

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = scenes.ASSETS
ICONS = os.path.join(ASSETS, "icons")
os.makedirs(ICONS, exist_ok=True)

# fonts
FD, FS, FE = D.F_DISPLAY, D.F_SERIF, D.F_EDITORIAL      # display serifs
FB, FU = D.F_BODY, D.F_UI                               # sans

W, H = D.SLIDE_W_IN, D.SLIDE_H_IN
MX, MY = D.MARGIN_X, D.MARGIN_Y
CW = D.CONTENT_W

prs = Presentation()
prs.slide_width = Emu(D.SLIDE_W_EMU)
prs.slide_height = Emu(D.SLIDE_H_EMU)


# =============================================================== asset caches ==
def icon(name: str, color: str = C.GOLD, lw: float = 3.0) -> str:
    """Render (and cache) one icon in a given colour."""
    key = f"{name}_{color}_{lw:.1f}.png"
    path = os.path.join(ICONS, key)
    if not os.path.exists(path):
        icons.render(name, px=300, color=color, lw=lw).save(path, optimize=True)
    return path


MAP_UK = (1800, 1500)
MAP_IN = (1500, 1500)


def build_maps(force: bool = False) -> None:
    def need(n):
        p = os.path.join(ASSETS, n)
        return force or not os.path.exists(p) or os.path.getsize(p) == 0

    if need("map_india.png"):
        mapgen.india_map(size=MAP_IN).save(os.path.join(ASSETS, "map_india.png"),
                                           optimize=True)
    layers = None
    keys = ["shadow", "glow_garhwal", "glow_kumaon", "garhwal", "kumaon",
            "districts", "rivers"]
    if any(need(f"map_uk_{k}.png") for k in keys):
        layers = mapgen.uttarakhand_layers(size=MAP_UK)
        for k in keys:
            layers[k].save(os.path.join(ASSETS, f"map_uk_{k}.png"), optimize=True)
    if need("map_uk_markers.png"):
        mapgen.kauthik_markers(size=MAP_UK).save(
            os.path.join(ASSETS, "map_uk_markers.png"), optimize=True)
    if need("map_uk_route.png"):
        mapgen.route_layer(size=MAP_UK).save(
            os.path.join(ASSETS, "map_uk_route.png"), optimize=True)
    if need("map_uk_full.png"):
        if layers is None:
            layers = mapgen.uttarakhand_layers(size=MAP_UK)
        from PIL import Image
        comp = Image.new("RGBA", MAP_UK, (0, 0, 0, 0))
        for k in ["shadow", "glow_garhwal", "glow_kumaon", "garhwal", "kumaon",
                  "districts", "rivers"]:
            comp.alpha_composite(layers[k])
        comp.alpha_composite(mapgen.kauthik_markers(size=MAP_UK))
        comp.save(os.path.join(ASSETS, "map_uk_full.png"), optimize=True)


# =================================================================== new slide ==
def new_slide(transition: str = "morph", dur: int = D.D_TRANSITION,
              bg: str = C.INK):
    slide = blank_slide(prs)
    set_slide_bg(slide, bg)
    tl = Timeline(slide)
    set_transition(slide, transition, dur)
    return slide, tl


# ================================================================ text helpers ==
def eyebrow(slide, x, y, txt, color=C.GOLD, w=9.0, align="l", size=D.T_EYEBROW,
            tick=True):
    """A tracked-out kicker. Left variant gets a small accent tick."""
    shapes = []
    if align == "l" and tick:
        shapes.append(rect(slide, x, y + 0.085, 0.34, 0.026, fill=color,
                           name="ebtick"))
        tx = x + 0.52
    else:
        tx = x
    shapes.append(text(slide, tx, y, w, 0.42,
                       [{"text": txt, "font": FU, "size": size, "color": color,
                         "track": D.TRACK_EYEBROW, "caps": True, "bold": True,
                         "align": align}], name="eyebrow"))
    return shapes


def flanked_eyebrow(slide, cy, txt, color=C.GOLD, half=1.05, gap=2.4):
    """Centre kicker with two short rules — used on dividers and closing."""
    cxc = W / 2
    tb = text(slide, cxc - 4.0, cy, 8.0, 0.4,
              [{"text": txt, "font": FU, "size": D.T_EYEBROW, "color": color,
                "track": D.TRACK_EYEBROW, "caps": True, "bold": True,
                "align": "c"}], name="eyebrow")
    l = rect(slide, cxc - gap - half, cy + 0.16, half, 0.02, fill=color, name="ebl")
    r = rect(slide, cxc + gap, cy + 0.16, half, 0.02, fill=color, name="ebr")
    return [l, r, tb]


def para(txt, font=FB, size=D.T_BODY, color=C.SNOW_DIM, bold=False,
         italic=False, align="l", track=0, line=1.28, caps=False,
         before=0.0, after=0.0, shadow=False, cs=False):
    d = dict(text=txt, font=font, size=size, color=color, bold=bold,
             italic=italic, align=align, line=line)
    if track:
        d["track"] = track
    if caps:
        d["caps"] = True
    if before:
        d["before"] = before
    if after:
        d["after"] = after
    if shadow:
        d["shadow"] = shadow
    if cs:
        d["cs"] = True
    return d


def rule(slide, x, y, w, color=C.GOLD, weight=1.1, alpha=None):
    return hairline(slide, x, y, w, color=color, weight=weight, alpha=alpha,
                    name="rule")


def connector(slide, x1, y1, x2, y2, color=C.GOLD_DEEP, weight=0.9, dash=None):
    cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     IN(x1), IN(y1), IN(x2), IN(y2))
    cxn.line.color.rgb = RGBColor(*D.rgb(color))
    cxn.line.width = Pt(weight)
    cxn.shadow.inherit = False
    if dash:
        from pptx.oxml.ns import qn
        ln = cxn.line._get_or_add_ln()
        d = ln.makeelement(qn("a:prstDash"), {"val": dash})
        ln.append(d)
    cxn._element._nvXxPr.cNvPr.set("name", "cxn")
    return cxn


# ===================================================================== cards ===
def feature_card(slide, x, y, w, h, icon_name, head, body, *, accent=C.GOLD,
                 head_size=D.T_H3, body_size=13.5, body_color=C.SNOW_DIM,
                 fill="0A1512", fill_alpha=0.40, meta=None, icon_lw=3.2,
                 body_top=1.28):
    """A translucent panel with an accent bar, icon, heading and body."""
    shapes = {}
    panel = rect(slide, x, y, w, h, fill=fill, alpha=fill_alpha,
                 line=C.GOLD_DEEP, line_w=0.75, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                 radius=0.055, name="card")
    shapes["card"] = panel
    shapes["bar"] = rect(slide, x, y, 0.055, h, fill=accent, name="cardbar")
    ic = picture(slide, icon(icon_name, accent, icon_lw), x + 0.30, y + 0.30,
                 0.72, 0.72, name="cardicon")
    shapes["icon"] = ic
    ty = y + 0.30
    tx = x + 1.24
    tw = w - 1.5
    paras = [para(head, font=FE, size=head_size, color=C.SNOW, bold=False,
                  line=1.02)]
    if meta:
        paras.append(para(meta, font=FU, size=10.5, color=accent, caps=True,
                          track=170, before=3, line=1.0))
    shapes["head"] = text(slide, tx, ty, tw, 1.0, paras, name="cardhead")
    shapes["body"] = text(slide, x + 0.30, y + body_top, w - 0.6,
                          h - body_top - 0.12,
                          [para(body, size=body_size, color=body_color,
                                line=1.30)], name="cardbody")
    return shapes


def animate_card(tl, shapes, delay, mode="rise"):
    getattr(tl, {"rise": "rise_in", "fade": "fade_in",
                 "drift": "drift_in"}[mode])(shapes["card"], delay=delay)
    tl.wipe_in(shapes["bar"], delay=delay + 120, dur=700, direction="up")
    tl.fade_in(shapes["icon"], delay=delay + 200)
    tl.fade_in(shapes["head"], delay=delay + 260)
    tl.fade_in(shapes["body"], delay=delay + 360)


def grid_positions(n, cols, x0, y0, cw, ch, gx, gy):
    return [(x0 + (i % cols) * (cw + gx), y0 + (i // cols) * (ch + gy))
            for i in range(n)]


# =============================================================== page furniture ==
def footer(slide, idx, label, mood_tone=C.GOLD):
    """Consistent baseline: index • section on the left, wordmark on the right."""
    text(slide, MX, H - 0.46, 6.0, 0.3,
         [para(f"{idx:02d}  ·  {label}", font=FU, size=9.5, color=C.SLATE_LIGHT,
               caps=True, track=190)], name="fnum")
    text(slide, W - MX - 4.0, H - 0.46, 4.0, 0.3,
         [para("KAUTHIK  ·  THE LIVING FAIRS OF UTTARAKHAND", font=FU, size=9.5,
               color=C.SLATE_LIGHT, caps=True, track=150, align="r")],
         name="fmark")


# ===================================================================== SLIDES ===
IDX = 0


def _idx():
    global IDX
    IDX += 1
    return IDX


# ---- 01 cover ---------------------------------------------------------------
def s_cover():
    slide, tl = new_slide(transition="none", bg=C.NIGHT)
    bg = scenes.scene(slide, tl, mood="dawn", rays="rays_top", mid=True,
                      ground="temple", ground_h=2.35, mist="mist_low",
                      particles="dust", scrim="bottom", flags=True,
                      mandala=(6.667, 3.05, 6.4), aipan="gold")
    # ghost wordmark in Devanagari, high and faint
    dev = text(slide, 0, 1.02, W, 0.9,
               [para("कौथिग", font=FS, size=40, color=C.GOLD_LIGHT, align="c",
                     track=220, cs=True)], name="devtitle")

    eb = flanked_eyebrow(slide, 2.62,
                         "A Cinematic Journey Through Himalayan Folk Culture")

    title = text(slide, 0, 2.98, W, 1.9,
                 [para("KAUTHIK", font=FD, size=D.T_HERO, color=C.SNOW,
                       align="c", track=60, bold=True,
                       shadow=dict(blur=34, dist=6, alpha=60))], name="herotitle")
    sub = text(slide, 0, 4.78, W, 0.8,
               [para("The Living Fairs of Uttarakhand", font=FS, size=32,
                     color=C.GOLD_LIGHT, align="c", italic=True,
                     shadow=dict(blur=20, dist=3, alpha=55))], name="herosub")
    r = rule(slide, W / 2 - 1.35, 5.66, 2.7, color=C.GOLD, weight=1.3)
    tag = text(slide, 0, 5.86, W, 0.4,
               [para("Garhwal   ·   Kumaon   ·   Devbhumi", font=FU, size=13,
                     color=C.SNOW_DIM, align="c", track=280, caps=True)],
               name="herotag")

    # ambient score — autoplays and loops as the living soundtrack
    aud = media.add_audio(slide, os.path.join(ASSETS, "score.m4a"),
                          x=W - 0.56, y=H - 0.56, size=0.40)

    tl.zoom_in(dev, delay=300, dur=1600)
    for s in eb:
        tl.fade_in(s, delay=700, dur=1200)
    tl.rise_in(title, delay=1050, dur=1700, rise=0.05)
    tl.rise_in(sub, delay=1750, dur=1500)
    tl.wipe_in(r, delay=2250, dur=900)
    tl.fade_in(tag, delay=2500, dur=1300)
    tl.play_media(aud, loop=True, vol=55000, kind="audio")
    tl.apply()
    notes(slide, K.NOTES["cover"])


# ---- 02 cold open (song) ----------------------------------------------------
def s_coldopen():
    slide, tl = new_slide(bg="0A0605")
    scenes.scene(slide, tl, mood="ember", rays="rays_top", mid=True,
                 ground="pines", ground_h=2.0, mist="mist_warm",
                 particles="embers", scrim="full", glow=(6.667, 3.2, 7.5, "gold"))
    eb = eyebrow(slide, MX, 1.15, "Before We Begin  ·  A Song From The Hills",
                 align="l")
    q = text(slide, MX, 2.15, W - 2 * MX, 2.6,
             [para("\u201cBedu pako baro maasa,", font=FE, size=46, color=C.SNOW,
                   italic=True, line=1.14, shadow=dict(blur=24, dist=4, alpha=55)),
              para("Narain\u2014 kaafal pako chaita\u2026\u201d", font=FE, size=46,
                   color=C.SNOW, italic=True, line=1.14,
                   shadow=dict(blur=24, dist=4, alpha=55))], name="quote")
    dev = text(slide, MX, 4.62, W - 2 * MX, 0.6,
               [para("बेडु पाको बारो मासा, नारैण\u2014 काफल पाको चैता",
                     font=FS, size=21, color=C.GOLD_LIGHT, cs=True, italic=True)],
               name="devline")
    tr = text(slide, MX, 5.42, 8.6, 1.2,
              [para("The wild fig ripens all twelve months \u2014 but the kafal, "
                    "the mountain bayberry, only in spring.", size=17,
                    color=C.SNOW_DIM, line=1.34),
               para("A song about the things that come but once a year. So is a "
                    "Kauthik.", size=17, color=C.GOLD_LIGHT, italic=True,
                    before=6, line=1.34, font=FS)], name="gloss")
    for s in eb:
        tl.fade_in(s, delay=300)
    tl.rise_in(q, delay=650, dur=1600, rise=0.04)
    tl.fade_in(dev, delay=1500, dur=1300)
    tl.fade_in(tr, delay=2000, dur=1300)
    footer(slide, _idx(), "Overture")
    tl.apply()
    notes(slide, K.NOTES["coldopen"])


# ---- 03 devbhumi ------------------------------------------------------------
def s_devbhumi():
    slide, tl = new_slide(bg="030B18")
    scenes.scene(slide, tl, mood="royal", rays="rays_cool", mid=True,
                 ground="pines_soft", ground_h=1.7, mist="mist_cold",
                 particles="dust", scrim="bottom", side="left")
    eb = eyebrow(slide, MX, 1.06, "The Setting  ·  Devbhumi, The Land Of The Gods")
    ttl = text(slide, MX, 1.62, 6.6, 2.0,
               [para("A young state with an", font=FE, size=40, color=C.SNOW,
                     line=1.06),
                para("ancient calendar of fairs", font=FE, size=40,
                     color=C.GOLD_LIGHT, line=1.06)], name="title")
    rule(slide, MX, 3.42, 2.2)
    body = text(slide, MX, 3.66, 5.7, 2.2,
                [para("Uttarakhand rises from the Gangetic plain to the high "
                      "Himalaya in one vertical sweep. Its people split into "
                      "two divisions with distinct tongues \u2014 Garhwal in "
                      "the west, Kumaon in the east.", size=15.5,
                      color=C.SNOW_DIM, line=1.32),
                 para("Formed in 2000, it is among India\u2019s youngest states "
                      "\u2014 and one of the last strongholds of a living fair "
                      "culture. Here, before the roads, the fair was "
                      "infrastructure.", size=15.5, color=C.SNOW_DIM, line=1.32,
                      before=8)], name="body")

    facts = [("13", "DISTRICTS"), ("2", "DIVISIONS"), ("2000", "STATEHOOD"),
             ("4", "CHAR DHAM")]
    fx = []
    for i, (n, l) in enumerate(facts):
        x = MX + i * 1.48
        a = text(slide, x, 6.12, 1.4, 0.55,
                 [para(n, font=FD, size=24, color=C.GOLD_LIGHT, bold=True)],
                 name="fn")
        b = text(slide, x, 6.60, 1.4, 0.35,
                 [para(l, font=FU, size=9.5, color=C.SLATE_LIGHT, caps=True,
                       track=120)], name="fl")
        fx += [a, b]

    # maps: India (highlight) with Uttarakhand inset
    mi = picture(slide, os.path.join(ASSETS, "map_india.png"), 7.55, 1.35,
                 5.1, 5.1, name="mapIndia")
    # UK inset, lower-right
    uk = picture(slide, os.path.join(ASSETS, "map_uk_full.png"), 9.15, 4.05,
                 3.55, 2.96, name="mapUK")
    cap = text(slide, 9.15, 6.95, 3.55, 0.3,
               [para("Uttarakhand  ·  Garhwal + Kumaon", font=FU, size=9,
                     color=C.SLATE_LIGHT, caps=True, track=120, align="c")],
               name="mapcap")

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.rise_in(ttl, delay=550, dur=1400)
    tl.fade_in(body, delay=1050, dur=1300)
    tl.zoom_in(mi, delay=800, dur=1600, start=0.9)
    tl.fade_in(uk, delay=1500, dur=1300)
    tl.fade_in(cap, delay=1900)
    tl.stagger(fx, start=1500, gap=140, mode="fade")
    footer(slide, _idx(), "The Land")
    tl.apply()
    notes(slide, K.NOTES["devbhumi"])


# ---- chapter divider --------------------------------------------------------
def chapter(num, name, sub, mood, *, particles="dust", ground=None,
            note_key=None, rays="rays_top"):
    slide, tl = new_slide(transition="smooth", dur=D.D_TRANSITION_SLOW,
                          bg=C.INK)
    scenes.scene(slide, tl, mood=mood, rays=rays, mid=True, ground=ground,
                 ground_h=2.0, mist="mist_low", particles=particles,
                 scrim="full", mandala=(10.4, 3.4, 6.0), aipan="gold")
    # ghost numeral
    numeral = text(slide, 6.7, 0.6, 6.2, 5.6,
                   [para(num, font=FD, size=250, color="4A4028", align="c",
                         bold=True)], name="numeral")
    eb = eyebrow(slide, MX, 3.0, f"Chapter {num}", color=C.GOLD)
    ttl = text(slide, MX, 3.5, 8.4, 1.5,
               [para(name, font=FD, size=64, color=C.SNOW, bold=True, track=30,
                     shadow=dict(blur=28, dist=5, alpha=55))], name="chtitle")
    rule(slide, MX + 0.02, 4.86, 2.4, weight=1.3)
    st = text(slide, MX, 5.06, 7.6, 0.8,
              [para(sub, font=FS, size=24, color=C.GOLD_LIGHT, italic=True)],
              name="chsub")
    tl.zoom_in(numeral, delay=200, dur=2000, start=0.86)
    for s in eb:
        tl.fade_in(s, delay=700)
    tl.rise_in(ttl, delay=1000, dur=1500, rise=0.05)
    tl.wipe_in(slide.shapes[-2] if False else st, delay=0) if False else None
    tl.fade_in(st, delay=1600, dur=1300)
    footer(slide, _idx(), name.title())
    tl.apply()
    if note_key:
        notes(slide, K.NOTES[note_key])


# ---- 05 the word ------------------------------------------------------------
def s_word():
    slide, tl = new_slide(bg="14130F")
    scenes.scene(slide, tl, mood="parchment", rays="rays_left", mid=False,
                 mist="mist_mid", particles="dust", scrim="veil",
                 texture="paper", mandala=(10.2, 3.8, 6.6), aipan="gold")
    eb = eyebrow(slide, MX, 0.95, "The Word  ·  Kauthig")
    big = text(slide, MX, 1.5, 11.5, 1.5,
               [para([{"text": "\u201cKauthig\u201d simply means ",
                       "font": FE, "size": 40, "color": C.SNOW},
                      {"text": "a fair", "font": FE, "size": 40,
                       "color": C.GOLD_LIGHT, "italic": True},
                      {"text": ".", "font": FE, "size": 40, "color": C.SNOW}],
                     line=1.1)], name="bigline")

    cards = [
        ("book", "A word in the mother tongue",
         "In Garhwali it is kauthig; in Kumaoni, kauthig or thol. Scholars trace "
         "it to the Sanskrit kautuka \u2014 curiosity, delight, spectacle. A "
         "reading, not a certainty \u2014 but a telling one."),
        ("hands", "Not an event \u2014 a gathering",
         "A Kauthik is less a scheduled programme than the day a scattered "
         "mountain people becomes, briefly, a single crowd. The word carries "
         "the joy of that assembly."),
        ("flag", "A festival given its own name",
         "Dehradun\u2019s fair began in 1998 as the \u2018Uttarakhand "
         "Mahotsav\u2019, then was renamed Kauthig \u2014 so the name itself "
         "would speak Garhwali. They gave a festival back its own word."),
    ]
    pos = grid_positions(3, 3, MX, 3.05, 3.62, 3.15, 0.24, 0)
    for i, ((ic, hd, bd), (x, y)) in enumerate(zip(cards, pos)):
        sh = feature_card(slide, x, y, 3.62, 3.15, ic, hd, bd,
                          head_size=20, body_size=13, body_top=1.78)
        animate_card(tl, sh, delay=900 + i * 320)

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.rise_in(big, delay=550, dur=1500)
    footer(slide, _idx(), "The Word")
    tl.apply()
    notes(slide, K.NOTES["word"])


# ---- 06 anatomy wheel -------------------------------------------------------
def s_anatomy():
    slide, tl = new_slide(bg="060D1A")
    scenes.scene(slide, tl, mood="night", rays="rays_cool", mid=True,
                 particles="fireflies", mist="mist_low", scrim="full")
    eb = eyebrow(slide, MX, 0.72, "The Institution  ·  Six Roles, One Day",
                 align="l")
    ttl = text(slide, MX, 0.98, 9.0, 0.8,
               [para("A Kauthik does six things at once", font=FE, size=30,
                     color=C.SNOW)], name="title")

    cx, cy = W / 2, 4.5
    rx, ry = 4.15, 1.9
    import math
    nodes = [(K.ANATOMY[i][0], K.ANATOMY[i][1]) for i in range(6)]
    short = ["anchored to a shrine", "the only marketplace", "struck on the dhol",
             "scattered kin meet", "millet, hemp, sweets", "songs find a stage"]
    # hub
    hub = ring(slide, cx, cy, 1.7, color=C.GOLD, weight=1.4, name="hub")
    hub2 = ring(slide, cx, cy, 1.42, color=C.GOLD_DEEP, weight=0.8, name="hub2")
    hublbl = text(slide, cx - 1.1, cy - 0.44, 2.2, 0.9,
                  [para("KAUTHIK", font=FD, size=22, color=C.GOLD_LIGHT,
                        align="c", bold=True, track=40),
                   para("one day, six roles", font=FS, size=13,
                        color=C.SNOW_DIM, align="c", italic=True, before=2)],
                  name="hublbl")

    node_shapes = []
    conns = []
    for i, (ic, hd) in enumerate(nodes):
        ang = -math.pi / 2 + i * (2 * math.pi / 6)
        nx, ny = cx + math.cos(ang) * rx, cy + math.sin(ang) * ry
        conns.append(connector(slide, cx + math.cos(ang) * 0.9,
                               cy + math.sin(ang) * 0.9,
                               nx - math.cos(ang) * 0.52,
                               ny - math.sin(ang) * 0.34,
                               color=C.GOLD_DEEP, weight=0.9))
        d = dot(slide, nx, ny, 0.92, fill="0C1A2E", line=C.GOLD, line_w=1.2,
                name="node")
        ig = picture(slide, icon(ic, C.GOLD_LIGHT, 3.4), nx - 0.31, ny - 0.31,
                     0.62, 0.62, name="nodeicon")
        right = math.cos(ang) >= -0.2
        top = math.sin(ang) < -0.5
        bottom = math.sin(ang) > 0.5
        if top:
            lx, ly, al = nx - 1.5, ny - 1.02, "c"
            lw_ = 3.0
        elif bottom:
            lx, ly, al = nx - 1.5, ny + 0.52, "c"
            lw_ = 3.0
        elif right:
            lx, ly, al = nx + 0.58, ny - 0.34, "l"
            lw_ = 2.6
        else:
            lx, ly, al = nx - 3.18, ny - 0.34, "r"
            lw_ = 2.6
        lbl = text(slide, lx, ly, lw_, 0.9,
                   [para(hd, font=FU, size=14.5, color=C.SNOW, bold=True,
                         caps=True, track=60, align=al),
                    para(short[i], font=FS, size=12.5, color=C.GOLD_LIGHT,
                         italic=True, align=al, before=1)], name="nodelbl")
        node_shapes.append((d, ig, lbl))

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    tl.zoom_in(hub, delay=800, dur=1300, start=0.6)
    tl.zoom_in(hub2, delay=850, dur=1300, start=0.6)
    tl.fade_in(hublbl, delay=1050)
    for i, (c, (d, ig, lbl)) in enumerate(zip(conns, node_shapes)):
        base = 1300 + i * 300
        tl.fade_in(c, delay=base, dur=700)
        tl.zoom_in(d, delay=base + 150, dur=800, start=0.4)
        tl.fade_in(ig, delay=base + 320)
        tl.fade_in(lbl, delay=base + 420)
    footer(slide, _idx(), "Anatomy Of A Fair")
    tl.apply()
    notes(slide, K.NOTES["anatomy"])


# ---- 08 timeline ------------------------------------------------------------
def s_timeline():
    slide, tl = new_slide(bg="030B18")
    scenes.scene(slide, tl, mood="royal", rays="rays_cool", mid=True,
                 ground="pines_soft", ground_h=1.5, mist="mist_cold",
                 particles="dust", scrim="full")
    eb = eyebrow(slide, MX, 0.66, "The History  ·  A Thousand Years Coming Down The Hill")
    ttl = text(slide, MX, 1.08, 10.5, 0.8,
               [para("How a fair became a movement, a market, and a memory",
                     font=FE, size=26, color=C.SNOW)], name="title")

    y0 = 4.4
    x0, x1 = MX + 0.1, W - MX - 0.1
    spine = rect(slide, x0, y0, x1 - x0, 0.028, fill=C.GOLD_DEEP, name="spine")
    items = K.TIMELINE
    n = len(items)
    step = (x1 - x0) / (n - 1)
    node_sh = []
    for i, (yr, hd, bd) in enumerate(items):
        x = x0 + i * step
        up = (i % 2 == 0)
        d = dot(slide, x, y0 + 0.014, 0.20, fill=C.GOLD_LIGHT, line=C.NIGHT,
                line_w=1.2, name="tnode")
        stem_y = (y0 - 0.62) if up else (y0 + 0.16)
        stem = vline(slide, x, min(stem_y, y0 + 0.014) if up else y0 + 0.16,
                     0.46, color=C.GOLD_DEEP, weight=0.8)
        if up:
            ty, yh, va = y0 - 2.28, 1.6, "b"
        else:
            ty, yh, va = y0 + 0.66, 1.7, "t"
        card_w = 2.42
        tx = min(max(x - card_w / 2, MX), W - MX - card_w)
        blk = text(slide, tx, ty, card_w, yh,
                   [para(yr, font=FU, size=13, color=C.GOLD_LIGHT, bold=True,
                         caps=True, track=80, align="c"),
                    para(hd, font=FE, size=16, color=C.SNOW, align="c",
                         before=2, line=1.02),
                    para(bd, font=FB, size=10.5, color=C.SNOW_DIM, align="c",
                         before=4, line=1.22)],
                   anchor=MSO_ANCHOR.BOTTOM if up else MSO_ANCHOR.TOP,
                   name="tcard")
        node_sh.append((d, stem, blk))

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    tl.wipe_in(spine, delay=800, dur=1800, direction="left")
    for i, (d, stem, blk) in enumerate(node_sh):
        base = 1100 + i * 340
        tl.zoom_in(d, delay=base, dur=600, start=0.3)
        tl.fade_in(stem, delay=base + 120, dur=500)
        tl.rise_in(blk, delay=base + 200, dur=900, rise=0.03)
    footer(slide, _idx(), "Timeline")
    tl.apply()
    notes(slide, K.NOTES["timeline"])


# ---- 09 1921 hero -----------------------------------------------------------
def s_1921():
    slide, tl = new_slide(bg="070E1C")
    scenes.scene(slide, tl, mood="dusk", rays="rays_right", mid=True,
                 ground="pines", ground_h=1.9, mist="mist_warm",
                 particles="dust", scrim="bottom_r", side="left",
                 glow=(10.6, 2.6, 6.2, "gold"))
    eb = eyebrow(slide, MX, 1.02, "The Hero Moment  ·  Bageshwar, 14 January 1921")
    yr = text(slide, MX, 1.30, 6.5, 1.5,
              [para("1921", font=FD, size=86, color=C.GOLD_LIGHT, bold=True,
                    shadow=dict(blur=30, dist=5, alpha=60))], name="year")
    hd = text(slide, MX, 3.06, 6.9, 1.1,
              [para("The fair that tore up", font=FE, size=32, color=C.SNOW,
                    line=1.05),
               para("forced labour", font=FE, size=32, color=C.GOLD_LIGHT,
                    line=1.05)], name="head")
    rule(slide, MX + 0.02, 4.10, 2.2)
    bd = text(slide, MX, 4.32, 6.6, 2.1,
              [para("At the Uttarayani Kauthik, on the sand where the Saryu "
                    "meets the Gomti, thousands of hill farmers ended "
                    "coolie-begar \u2014 forced, unpaid labour. Led by Badri "
                    "Datt Pandey and Hargovind Pant, they tore the registers "
                    "and threw them into the river.", size=15.5, color=C.SNOW_DIM,
                    line=1.32),
               para("Nothing else could gather thousands from a hundred "
                    "valleys on a known day. The fair was the only public "
                    "square there was.", size=15.5, color=C.GOLD_LIGHT,
                    italic=True, font=FS, before=8, line=1.3)], name="body")

    # right pull-quote panel
    panel = rect(slide, 8.55, 2.5, 3.95, 2.7, fill="10151F", alpha=0.5,
                 line=C.GOLD_DEEP, line_w=0.8, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
                 radius=0.05, name="qpanel")
    pq = text(slide, 8.85, 2.78, 3.4, 2.2,
              [para("\u201cKumaon Kesari\u201d", font=FE, size=22,
                    color=C.SNOW, italic=True),
               para("Badri Datt Pandey \u2014 the Lion of Kumaon", size=12,
                    color=C.SNOW_DIM, before=3, line=1.25),
               para("\u201cJannayak\u201d", font=FE, size=22, color=C.SNOW,
                    italic=True, before=12),
               para("Hargovind Pant \u2014 leader of the people", size=12,
                    color=C.SNOW_DIM, before=3, line=1.25)], name="pquote")
    note = text(slide, 8.55, 5.42, 3.95, 0.7,
                [para("Every January, the gesture is still repeated on that "
                      "same sandbank.", size=12, color=C.GOLD_LIGHT, italic=True,
                      align="c", font=FS, line=1.3)], name="qnote")

    for s in eb:
        tl.fade_in(s, delay=300)
    tl.rise_in(yr, delay=650, dur=1700, rise=0.04)
    tl.rise_in(hd, delay=1300, dur=1300)
    tl.fade_in(bd, delay=1750, dur=1300)
    tl.drift_in(panel, delay=1500, dur=1200, dx=0.03)
    tl.fade_in(pq, delay=1900)
    tl.fade_in(note, delay=2300)
    footer(slide, _idx(), "1921  ·  Coolie-Begar")
    tl.apply()
    notes(slide, K.NOTES["1921"])


# ---- 11 map -----------------------------------------------------------------
def s_map():
    slide, tl = new_slide(bg="0C1826")
    scenes.scene(slide, tl, mood="snow", rays="rays_cool", mid=True,
                 ground="pines_soft", ground_h=1.4, mist="mist_cold",
                 particles="snowfall", scrim="bottom", side="left")
    eb = eyebrow(slide, MX, 0.78, "The Land  ·  Where The Fairs Stand")
    ttl = text(slide, MX, 1.24, 5.6, 1.3,
               [para("A map you could", font=FE, size=34, color=C.SNOW,
                     line=1.05),
                para("walk by festival", font=FE, size=34, color=C.GOLD_LIGHT,
                     line=1.05)], name="title")
    rule(slide, MX + 0.02, 2.66, 2.0)

    sites = [
        ("BAGESHWAR", "Uttarayani Kauthik", "Saryu\u2013Gomti confluence; once "
         "the largest fair in Kumaon."),
        ("JAULJIBI", "Indo-Nepal trade mela", "Where the Kali meets the Gori "
         "\u2014 a bazaar in two countries."),
        ("DEVIDHURA", "Bagwal, at Barahi Devi", "Four clans meet in ritual "
         "combat on Raksha Bandhan."),
        ("NAUTI", "Nanda Devi Raj Jat", "Start of the 280 km pilgrimage, once "
         "every twelve years."),
    ]
    sh = []
    for i, (nm, sub, bd) in enumerate(sites):
        y = 2.94 + i * 0.98
        d = dot(slide, MX + 0.09, y + 0.12, 0.16, fill=C.GOLD_LIGHT, name="mdot")
        t = text(slide, MX + 0.42, y, 5.7, 0.95,
                 [para([{"text": nm + "   ", "font": FU, "size": 14.5,
                         "color": C.SNOW, "bold": True, "caps": True,
                         "track": 60},
                        {"text": sub, "font": FS, "size": 14.5,
                         "color": C.GOLD_LIGHT, "italic": True}]),
                  para(bd, size=11.5, color=C.SNOW_DIM, before=2, line=1.2)],
                 name="msite")
        sh.append((d, t))

    # legend
    lg = text(slide, MX, 6.86, 6.0, 0.3,
              [para([{"text": "\u25CF ", "font": FU, "size": 11,
                      "color": C.PINE},
                     {"text": "Garhwal        ", "font": FU, "size": 10,
                      "color": C.SNOW_DIM, "caps": True, "track": 100},
                     {"text": "\u25CF ", "font": FU, "size": 11,
                      "color": "3E7AB0"},
                     {"text": "Kumaon", "font": FU, "size": 10,
                      "color": C.SNOW_DIM, "caps": True, "track": 100}])],
              name="legend")

    # map layers (registered, progressive reveal)
    mapx, mapy, mapw = 6.35, 1.15, 6.35
    maph = mapw * MAP_UK[1] / MAP_UK[0]
    L = {}
    for k in ["shadow", "glow_garhwal", "glow_kumaon", "garhwal", "kumaon",
              "districts", "rivers", "markers"]:
        L[k] = picture(slide, os.path.join(ASSETS, f"map_uk_{k}.png"),
                       mapx, mapy, mapw, maph, name=f"map_{k}")

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.rise_in(ttl, delay=550, dur=1300)
    tl.fade_in(L["shadow"], delay=700, dur=900)
    tl.wipe_in(L["garhwal"], delay=1000, dur=1100, direction="left")
    tl.fade_in(L["glow_garhwal"], delay=1000, dur=1200)
    tl.wipe_in(L["kumaon"], delay=1350, dur=1100, direction="right")
    tl.fade_in(L["glow_kumaon"], delay=1350, dur=1200)
    tl.fade_in(L["districts"], delay=1750, dur=900)
    tl.fade_in(L["rivers"], delay=2000, dur=1100)
    tl.fade_in(lg, delay=1900)
    tl.fade_in(L["markers"], delay=2300, dur=1200)
    for i, (d, t) in enumerate(sh):
        tl.zoom_in(d, delay=2500 + i * 240, dur=600, start=0.3)
        tl.drift_in(t, delay=2560 + i * 240, dur=900, dx=-0.02)
    footer(slide, _idx(), "The Map")
    tl.apply()
    notes(slide, K.NOTES["map"])


# ---- 12 calendar wheel ------------------------------------------------------
def s_calendar():
    slide, tl = new_slide(bg="03110C")
    scenes.scene(slide, tl, mood="forest", rays="rays_top", mid=True,
                 particles="pollen", mist="mist_low", scrim="full")
    eb = eyebrow(slide, MX, 0.72, "The Rhythm  ·  There Is No Single Kauthik",
                 align="l")
    ttl = text(slide, MX, 1.16, 9.5, 0.8,
               [para("A whole calendar, pinned to sun and seed", font=FE,
                     size=30, color=C.SNOW)], name="title")

    import math
    cx, cy, R = W / 2, 4.42, 2.42
    Rx, Ry = 4.3, 2.42
    ring(slide, cx, cy, R * 2, color=C.GOLD_DEEP, weight=0.8, name="ring")
    hub = text(slide, cx - 1.2, cy - 0.5, 2.4, 1.0,
               [para("THE YEAR", font=FD, size=20, color=C.GOLD_LIGHT,
                     align="c", bold=True, track=60),
                para("Chaitra \u2192 Makar", font=FS, size=12.5, color=C.SNOW_DIM,
                     align="c", italic=True, before=2)], name="hub")
    node_sh = []
    for i, (mon, date, fest, desc) in enumerate(K.CALENDAR):
        ang = -math.pi / 2 + i * (2 * math.pi / 8)
        nx, ny = cx + math.cos(ang) * Rx, cy + math.sin(ang) * Ry
        d = dot(slide, nx, ny, 0.30, fill=C.GOLD_LIGHT, line=C.FOREST_DEEP,
                line_w=1.0, name="cnode")
        right = math.cos(ang) >= -0.15
        top = math.sin(ang) < -0.55
        bot = math.sin(ang) > 0.55
        if top:
            lx, ly, al, wd = nx - 1.4, ny - 1.06, "c", 2.8
        elif bot:
            lx, ly, al, wd = nx - 1.4, ny + 0.34, "c", 2.8
        elif right:
            lx, ly, al, wd = nx + 0.34, ny - 0.42, "l", 2.9
        else:
            lx, ly, al, wd = nx - 3.24, ny - 0.42, "r", 2.9
        lbl = text(slide, lx, ly, wd, 1.0,
                   [para([{"text": mon + "  ", "font": FU, "size": 12.5,
                           "color": C.SNOW, "bold": True, "caps": True,
                           "track": 60},
                          {"text": date, "font": FU, "size": 9.5,
                           "color": C.SLATE_LIGHT, "caps": True}], align=al),
                    para(fest, font=FS, size=13, color=C.GOLD_LIGHT,
                         italic=True, align=al, before=1, line=1.05)],
                   name="clbl")
        node_sh.append((d, lbl))

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    tl.fade_in(hub, delay=900)
    for i, (d, lbl) in enumerate(node_sh):
        base = 1100 + i * 230
        tl.zoom_in(d, delay=base, dur=500, start=0.3)
        tl.fade_in(lbl, delay=base + 120)
    footer(slide, _idx(), "Wheel Of The Year")
    tl.apply()
    notes(slide, K.NOTES["calendar"])


# ---- 14 music ---------------------------------------------------------------
def s_music():
    slide, tl = new_slide(bg="070505")
    scenes.scene(slide, tl, mood="ember", rays="rays_top", mid=True,
                 ground="pines", ground_h=1.7, mist="mist_warm",
                 particles="embers", scrim="full")
    eb = eyebrow(slide, MX, 0.72, "The Sound  ·  It Begins When The Dhol Is Struck")
    ttl = text(slide, MX, 1.16, 10.0, 0.8,
               [para("Drum, horn, and the circle of dancers", font=FE, size=30,
                     color=C.SNOW)], name="title")

    # instruments row (icon grid)
    lab = text(slide, MX, 2.12, 6.0, 0.3,
               [para("The Instruments", font=FU, size=12, color=C.GOLD_LIGHT,
                     caps=True, track=180)], name="ilab")
    inst_sh = []
    for i, (ic, nm, ds) in enumerate(K.INSTRUMENTS):
        x = MX + (i % 3) * 2.28
        y = 2.52 + (i // 3) * 2.0
        d = dot(slide, x + 0.55, y + 0.55, 1.06, fill="1A0C0A", line=C.GOLD_DEEP,
                line_w=1.0, name="idot")
        ig = picture(slide, icon(ic, C.GOLD_LIGHT, 3.2), x + 0.19, y + 0.19,
                     0.72, 0.72, name="iicon")
        t = text(slide, x - 0.05, y + 1.16, 2.15, 0.9,
                 [para(nm, font=FU, size=12.5, color=C.SNOW, bold=True,
                       caps=True, track=70, align="c"),
                  para(ds, size=10, color=C.SNOW_DIM, align="c", before=1,
                       line=1.16)], name="ilbl")
        inst_sh.append((d, ig, t))

    # song forms column
    vln = vline(slide, 7.55, 2.2, 4.4, color=C.GOLD_DEEP, weight=0.8)
    slab = text(slide, 7.85, 2.12, 5.0, 0.3,
                [para("The Songs", font=FU, size=12, color=C.GOLD_LIGHT,
                      caps=True, track=180)], name="slab")
    song_sh = []
    for i, (nm, ds) in enumerate(K.SONG_FORMS):
        y = 2.56 + i * 0.72
        t = text(slide, 7.85, y, 5.0, 0.7,
                 [para([{"text": nm + "   ", "font": FU, "size": 13.5,
                         "color": C.GOLD_LIGHT, "bold": True, "caps": True,
                         "track": 80},
                        {"text": ds, "font": FB, "size": 12, "color": C.SNOW_DIM}],
                       line=1.18)], name="song")
        song_sh.append(t)

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    tl.fade_in(lab, delay=800)
    for i, (d, ig, t) in enumerate(inst_sh):
        base = 950 + i * 200
        tl.zoom_in(d, delay=base, dur=600, start=0.4)
        tl.fade_in(ig, delay=base + 150)
        tl.fade_in(t, delay=base + 230)
    tl.wipe_in(vln, delay=1100, dur=1000, direction="up")
    tl.fade_in(slab, delay=1300)
    tl.stagger(song_sh, start=1500, gap=200, mode="drift", dx=0.02)
    footer(slide, _idx(), "The Sound")
    tl.apply()
    notes(slide, K.NOTES["music"])


# ---- 15 dance ---------------------------------------------------------------
def s_dance():
    slide, tl = new_slide(bg="070E1C")
    scenes.scene(slide, tl, mood="dusk", rays="rays_right", mid=True,
                 ground="stalls", ground_h=1.7, mist="mist_warm",
                 particles="dust", scrim="full")
    eb = eyebrow(slide, MX, 0.72, "The Dance  ·  A Circle, Not A Stage")
    ttl = text(slide, MX, 1.16, 10.5, 0.8,
               [para("Steps with a job to do", font=FE, size=30, color=C.SNOW)],
               name="title")
    pos = grid_positions(4, 2, MX, 2.2, 5.5, 2.28, 0.35, 0.24)
    for i, ((ic, hd, meta, bd), (x, y)) in enumerate(zip(K.DANCES, pos)):
        sh = feature_card(slide, x, y, 5.5, 2.28, ic, hd, bd, meta=meta,
                          head_size=21, body_size=12)
        animate_card(tl, sh, delay=850 + i * 280)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    footer(slide, _idx(), "Folk Dance")
    tl.apply()
    notes(slide, K.NOTES["dance"])


# ---- 16 dance film (living background video) --------------------------------
def s_dance_film():
    """A full-bleed cinematic interlude: an original animated dance loop that
    plays and loops on its own behind a minimal caption."""
    slide, tl = new_slide(transition="smooth", dur=D.D_TRANSITION_SLOW, bg=C.INK)
    vid = media.add_video(slide, os.path.join(ASSETS, "dance.mp4"),
                          os.path.join(ASSETS, "dance_poster.jpg"),
                          -0.06, -0.05, W + 0.12, H + 0.10)
    # readability scrim over the moving footage
    picture(slide, os.path.join(ASSETS, "scrim_bottom.png"), 0, 0, W, H,
            name="scrim")
    eb = eyebrow(slide, MX, 0.72, "The Dance  ·  The Fair Remembering Itself",
                 color=C.GOLD_LIGHT)
    ln = text(slide, MX, H - 2.72, W - 2 * MX, 1.7,
              [para("When the dhol strikes, the circle forms itself \u2014",
                    font=FE, size=34, color=C.SNOW, italic=True, line=1.14,
                    shadow=dict(blur=22, dist=4, alpha=60)),
               para("no stage, no front row. Everyone is in it.",
                    font=FE, size=34, color=C.GOLD_LIGHT, italic=True, line=1.14,
                    shadow=dict(blur=22, dist=4, alpha=60))], name="filmline")
    cap = text(slide, MX, H - 0.94, W - 2 * MX, 0.4,
               [para("Chholiya  ·  Jhora  ·  Chhapeli  ·  Jhumeila  ·  Thadya",
                     font=FU, size=12.5, color=C.SNOW_DIM, caps=True, track=220)],
               name="filmcap")
    tl.play_media(vid, loop=True, kind="video")
    for s in eb:
        tl.fade_in(s, delay=400, dur=1400)
    tl.rise_in(ln, delay=900, dur=1700, rise=0.03)
    tl.fade_in(cap, delay=1900, dur=1400)
    footer(slide, _idx(), "The Dance, Alive")
    tl.apply()
    notes(slide, K.NOTES.get("dance_film", K.NOTES["dance"]))


# ---- 17 attire --------------------------------------------------------------
def s_attire():
    slide, tl = new_slide(bg="14130F")
    scenes.scene(slide, tl, mood="parchment", rays="rays_left", mid=False,
                 mist="mist_mid", particles="dust", scrim="veil",
                 texture="wool", mandala=(10.4, 3.7, 6.2), aipan="gold")
    eb = eyebrow(slide, MX, 0.9, "The Cloth  ·  Worn Wealth, Worn Memory")
    ttl = text(slide, MX, 1.4, 11.0, 0.9,
               [para("What a Kauthik wears", font=FE, size=34, color=C.SNOW)],
               name="title")
    pos = grid_positions(3, 3, MX, 2.75, 3.62, 3.5, 0.24, 0)
    for i, ((ic, hd, bd), (x, y)) in enumerate(zip(K.ATTIRE, pos)):
        sh = feature_card(slide, x, y, 3.62, 3.5, ic, hd, bd, head_size=19,
                          body_size=12.5, meta=None)
        animate_card(tl, sh, delay=850 + i * 300)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.rise_in(ttl, delay=550, dur=1300)
    footer(slide, _idx(), "Traditional Attire")
    tl.apply()
    notes(slide, K.NOTES["attire"])


# ---- 18 crafts --------------------------------------------------------------
def s_crafts():
    slide, tl = new_slide(bg="03110C")
    scenes.scene(slide, tl, mood="forest", rays="rays_top", mid=True,
                 ground="pines_soft", ground_h=1.4, mist="mist_low",
                 particles="pollen", scrim="full")
    eb = eyebrow(slide, MX, 0.72, "The Hands  ·  The Bazaar That Is Also A Museum")
    ttl = text(slide, MX, 1.16, 11.0, 0.8,
               [para("Made by hand, now protected by law", font=FE, size=29,
                     color=C.SNOW)], name="title")
    pos = grid_positions(6, 3, MX, 2.14, 3.62, 2.28, 0.24, 0.22)
    for i, (row, (x, y)) in enumerate(zip(K.CRAFTS, pos)):
        ic, hd = row[0], row[1]
        if len(row) == 4:
            meta, bd = row[2], row[3]
        else:
            meta, bd = None, row[2]
        accent = C.GOLD if i < 5 else C.COPPER_LIGHT
        sh = feature_card(slide, x, y, 3.62, 2.28, ic, hd, bd, meta=meta,
                          head_size=17, body_size=11, accent=accent)
        animate_card(tl, sh, delay=800 + i * 220)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    footer(slide, _idx(), "Local Handicrafts")
    tl.apply()
    notes(slide, K.NOTES["crafts"])


# ---- 19 food ----------------------------------------------------------------
def s_food():
    slide, tl = new_slide(bg="070505")
    scenes.scene(slide, tl, mood="ember", rays="rays_top", mid=True,
                 ground="stalls", ground_h=1.7, mist="mist_warm",
                 particles="embers", scrim="full")
    eb = eyebrow(slide, MX, 0.72, "The Meal  ·  Millet, Hemp, And A Sweet You Wear")
    ttl = text(slide, MX, 1.16, 11.0, 0.8,
               [para("The taste people name first when they leave", font=FE,
                     size=28, color=C.SNOW)], name="title")

    def col(x, head_ic, head, items, accent):
        vline(slide, x - 0.32, 2.34, 4.3, color=C.GOLD_DEEP, weight=0.8) \
            if x > 7 else None
        picture(slide, icon(head_ic, accent, 3.2), x, 2.16, 0.5, 0.5,
                name="foodhead")
        hh = text(slide, x + 0.66, 2.2, 5.0, 0.4,
                  [para(head, font=FU, size=13, color=accent, caps=True,
                        track=160, bold=True)], name="foodlab")
        sh = [hh]
        for i, (nm, ds) in enumerate(items):
            y = 2.92 + i * 0.94
            t = text(slide, x, y, 5.2, 0.9,
                     [para(nm, font=FE, size=17, color=C.SNOW, line=1.0),
                      para(ds, size=11.5, color=C.SNOW_DIM, before=2, line=1.2)],
                     name="fooditem")
            sh.append(t)
        return sh

    left = col(MX, "pot", "Savoury  ·  From The Earth", K.SAVOURY, C.GOLD_LIGHT)
    right = col(7.35, "sweet", "Sweet  ·  From The Festival", K.SWEETS,
                C.SAFFRON_LIGHT)

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    tl.stagger(left, start=850, gap=190, mode="drift", dx=-0.02)
    tl.stagger(right, start=1100, gap=190, mode="drift", dx=0.02)
    footer(slide, _idx(), "Cuisine")
    tl.apply()
    notes(slide, K.NOTES["food"])


# ---- 21 spirit --------------------------------------------------------------
def s_spirit():
    slide, tl = new_slide(bg="030B18")
    scenes.scene(slide, tl, mood="royal", rays="rays_cool", mid=True,
                 ground="temple", ground_h=2.5, mist="mist_cold",
                 particles="fireflies", scrim="bottom", side="left",
                 glow=(6.3, 3.0, 5.5, "cool"))
    eb = eyebrow(slide, MX, 0.86, "The Spirit  ·  Goddess, Oath And Ordeal")
    ttl = text(slide, MX, 1.32, 6.6, 1.4,
               [para("A pilgrimage that is", font=FE, size=33, color=C.SNOW,
                     line=1.05),
                para("really a farewell", font=FE, size=33, color=C.GOLD_LIGHT,
                     line=1.05)], name="title")
    rule(slide, MX + 0.02, 2.78, 2.0)
    bd = text(slide, MX, 3.0, 5.9, 2.1,
              [para("The Nanda Devi Raj Jat \u2014 the \u2018Himalayan "
                    "Mahakumbh\u2019 \u2014 walks the goddess from her "
                    "birth-village to her husband\u2019s home across the high "
                    "Himalaya. A four-horned ram, born for the task, leads the "
                    "way and is released at Homkund.", size=15.5,
                    color=C.SNOW_DIM, line=1.34),
               para("At Bagwal, four clans once drew blood for the goddess. "
                    "Since 2013, fruit and flowers stand in for the stones \u2014 "
                    "a tradition amending itself from the inside.", size=15.5,
                    color=C.GOLD_LIGHT, italic=True, font=FS, before=8,
                    line=1.32)], name="body")

    # Raj Jat stat strip
    sh = []
    for i, (n, u, d) in enumerate(K.RAJ_JAT):
        x = 7.2 + (i % 2) * 2.95
        y = 2.4 + (i // 2) * 1.72
        big = text(slide, x, y, 2.7, 0.9,
                   [para(n, font=FD, size=52, color=C.GOLD_LIGHT, bold=True)],
                   name="rjn")
        u_ = text(slide, x + 0.04, y + 0.92, 2.7, 0.3,
                  [para(u, font=FU, size=11, color=C.SNOW, caps=True, track=140,
                        bold=True)], name="rju")
        d_ = text(slide, x + 0.04, y + 1.2, 2.75, 0.5,
                  [para(d, size=10.5, color=C.SNOW_DIM, line=1.16)], name="rjd")
        sh += [(big, u_, d_)]

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.rise_in(ttl, delay=550, dur=1300)
    tl.fade_in(bd, delay=1050, dur=1300)
    for i, (a, b, c) in enumerate(sh):
        base = 1200 + i * 260
        tl.rise_in(a, delay=base, dur=1000, rise=0.04)
        tl.fade_in(b, delay=base + 200)
        tl.fade_in(c, delay=base + 300)
    footer(slide, _idx(), "Spiritual Significance")
    tl.apply()
    notes(slide, K.NOTES["spirit"])


# ---- 23 economy -------------------------------------------------------------
def s_economy():
    slide, tl = new_slide(bg="14130F")
    scenes.scene(slide, tl, mood="parchment", rays="rays_left", mid=True,
                 ground="stalls", ground_h=1.6, mist="mist_mid",
                 particles="dust", scrim="full", texture="paper")
    eb = eyebrow(slide, MX, 0.72, "The Marketplace  ·  What The Fair Is Worth")
    ttl = text(slide, MX, 1.16, 11.5, 0.9,
               [para([{"text": "The passes closed in ", "font": FE, "size": 30,
                       "color": C.SNOW},
                      {"text": "1962", "font": FE, "size": 30,
                       "color": C.GOLD_LIGHT},
                      {"text": ". The fair changed its trade.", "font": FE,
                       "size": 30, "color": C.SNOW}])], name="title")

    def column(x, head, items, accent, icon_name):
        picture(slide, icon(icon_name, accent, 3.0), x, 2.28, 0.52, 0.52,
                name="ecicon")
        hh = text(slide, x + 0.68, 2.32, 4.6, 0.4,
                  [para(head, font=FU, size=13, color=accent, caps=True,
                        track=150, bold=True)], name="echead")
        sh = [hh]
        for i, it in enumerate(items):
            y = 3.02 + i * 0.62
            d = dot(slide, x + 0.09, y + 0.11, 0.12, fill=accent, name="ecdot")
            t = text(slide, x + 0.34, y, 4.7, 0.5,
                     [para(it, size=14, color=C.SNOW_DIM, line=1.15)],
                     name="ecitem")
            sh += [d, t]
        return sh

    left = column(MX, "Then  ·  Across The Passes", K.TRADE_THEN, C.COPPER_LIGHT,
                  "cart")
    # divider arrow
    arr = text(slide, 6.35, 3.7, 0.7, 0.8,
               [para("\u2192", font=FB, size=40, color=C.GOLD, align="c")],
               name="arrow")
    right = column(7.35, "Now  ·  Named And Protected", K.TRADE_NOW,
                   C.GOLD_LIGHT, "gi_seal")

    cap = text(slide, MX, 6.5, 11.0, 0.6,
               [para("The fair stopped being where goods crossed a border, and "
                     "became where goods acquire an identity.", font=FS,
                     size=16, color=C.GOLD_LIGHT, italic=True, align="c",
                     line=1.25)], name="eccap")

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    tl.stagger(left, start=850, gap=150, mode="drift", dx=-0.02)
    tl.fade_in(arr, delay=1500, dur=900)
    tl.stagger(right, start=1650, gap=150, mode="drift", dx=0.02)
    tl.fade_in(cap, delay=2700)
    footer(slide, _idx(), "Economic Contribution")
    tl.apply()
    notes(slide, K.NOTES["economy"])


# ---- 24 tourism (chart + counter) -------------------------------------------
def s_tourism():
    slide, tl = new_slide(bg="0C1826")
    scenes.scene(slide, tl, mood="snow", rays="rays_cool", mid=True,
                 ground="pines_soft", ground_h=1.4, mist="mist_cold",
                 particles="snowfall", scrim="bottom", side="left")
    eb = eyebrow(slide, MX, 0.82, "Tourism  ·  A Counter-Weight To The Crush")
    ttl = text(slide, MX, 1.28, 6.4, 1.3,
               [para("The fairs spread the", font=FE, size=32, color=C.SNOW,
                     line=1.05),
                para("crowd across the year", font=FE, size=32,
                     color=C.GOLD_LIGHT, line=1.05)], name="title")
    rule(slide, MX + 0.02, 2.72, 2.0)
    bd = text(slide, MX, 2.94, 5.7, 2.2,
              [para("Char Dham pilgrimage is the state\u2019s engine \u2014 but it "
                    "peaks in a few months, on a few fragile routes.", size=15,
                    color=C.SNOW_DIM, line=1.32),
               para("The Kauthik calendar already exists, is distributed across "
                    "the map, and needs no new construction. It is the "
                    "off-season made sellable.", size=15, color=C.GOLD_LIGHT,
                    italic=True, font=FS, before=8, line=1.3)], name="body")

    # counting hero stat
    cnt = ["18.4", "31.7", "42.9", "46.7", "51.1"]
    stat_shapes = []
    for v in cnt:
        t = text(slide, MX, 5.3, 4.0, 1.0,
                 [para(v, font=FD, size=58, color=C.GOLD_LIGHT, bold=True)],
                 name="statcnt")
        stat_shapes.append(t)
    unit = text(slide, MX + 0.05, 6.42, 6.0, 0.4,
                [para("LAKH CHAR DHAM PILGRIMS  ·  2025", font=FU, size=11,
                      color=C.SNOW, caps=True, track=140, bold=True)],
                name="statunit")

    # bar chart, right
    bars = [("2021", 0.36, "19.6"), ("2022", 0.86, "46.3"), ("2023", 0.98, "56.1"),
            ("2024", 0.82, "46.7"), ("2025", 0.90, "51.1")]
    bx0, by0, bw, bh_max = 7.5, 6.2, 0.86, 3.2
    axis = rect(slide, bx0 - 0.2, by0 + 0.02, 5.0, 0.02, fill=C.SLATE, name="axis")
    bar_sh = []
    for i, (yr, frac, val) in enumerate(bars):
        x = bx0 + i * (bw + 0.22)
        bh = bh_max * frac
        highlight = yr in ("2025",)
        col = C.GOLD_LIGHT if highlight else D.mix_hex(C.HAZE_BLUE, C.SNOW, 0.15)
        b = rect(slide, x, by0 - bh, bw, bh, fill=col,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08, name="bar")
        vlab = text(slide, x - 0.15, by0 - bh - 0.42, bw + 0.3, 0.35,
                    [para(val, font=FU, size=12, color=C.SNOW, bold=True,
                          align="c")], name="barval")
        ylab = text(slide, x - 0.15, by0 + 0.08, bw + 0.3, 0.3,
                    [para(yr, font=FU, size=10.5, color=C.SLATE_LIGHT,
                          align="c", caps=True, track=60)], name="baryr")
        bar_sh.append((b, vlab, ylab))
    chart_cap = text(slide, 7.3, 1.5, 5.2, 0.7,
                     [para("Char Dham pilgrims, in lakh \u2014 the annual "
                           "spike the fairs can help flatten.", font=FU,
                           size=11, color=C.SLATE_LIGHT, align="r", line=1.3,
                           caps=True, track=80)], name="chartcap")

    for s in eb:
        tl.fade_in(s, delay=250)
    tl.rise_in(ttl, delay=550, dur=1300)
    tl.fade_in(bd, delay=1050, dur=1300)
    tl.fade_in(chart_cap, delay=900)
    tl.wipe_in(axis, delay=1000, dur=900, direction="left")
    for i, (b, vlab, ylab) in enumerate(bar_sh):
        base = 1200 + i * 260
        tl.wipe_in(b, delay=base, dur=900, direction="up")
        tl.fade_in(vlab, delay=base + 500)
        tl.fade_in(ylab, delay=base + 200)
    tl.count_up(stat_shapes, delay=1500, step=170)
    tl.fade_in(unit, delay=1500 + len(cnt) * 170)
    footer(slide, _idx(), "Tourism Impact")
    tl.apply()
    notes(slide, K.NOTES["tourism"])


# ---- 26 modern --------------------------------------------------------------
def s_modern():
    slide, tl = new_slide(bg="04101F")
    scenes.scene(slide, tl, mood="dawn", rays="rays_top", mid=True,
                 ground="stalls", ground_h=1.8, mist="mist_low",
                 particles="dust", scrim="full", flags=True)
    eb = eyebrow(slide, MX, 0.72, "The Modern Era  ·  The Fair Followed Its People")
    ttl = text(slide, MX, 1.16, 11.0, 0.8,
               [para("Down the mountain, onto the screen", font=FE, size=30,
                     color=C.SNOW)], name="title")
    cards = [
        ("homestay", "The city Kauthig",
         "Dehradun\u2019s Gad Kauthig marks statehood each November; Lucknow, "
         "Delhi and Mumbai run their own \u2014 migrant associations keeping the "
         "format, and the community, alive far from the hills."),
        ("signal", "A dialect goes viral",
         "Garhwali and Kumaoni folk music has found a real audience on "
         "short-video platforms for the first time \u2014 old songs, new reach."),
        ("hands", "Who the fair pays",
         "Women\u2019s self-help groups and collectives now run a large share "
         "of the craft and food stalls, changing who earns from the crowd."),
        ("seedling", "A tradition editing itself",
         "Bagwal swapping stones for fruit shows a living custom reforming from "
         "the inside \u2014 not preserved in amber, but still deciding."),
    ]
    pos = grid_positions(4, 2, MX, 2.2, 5.5, 2.24, 0.35, 0.22)
    for i, ((ic, hd, bd), (x, y)) in enumerate(zip(cards, pos)):
        sh = feature_card(slide, x, y, 5.5, 2.24, ic, hd, bd, head_size=20,
                          body_size=12)
        animate_card(tl, sh, delay=850 + i * 280)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    footer(slide, _idx(), "Modern Transformation")
    tl.apply()
    notes(slide, K.NOTES["modern"])


# ---- 27 challenges ----------------------------------------------------------
def s_challenges():
    slide, tl = new_slide(bg="060D1A")
    scenes.scene(slide, tl, mood="night", rays="rays_cool", mid=True,
                 ground="pines_soft", ground_h=1.6, mist="mist_low",
                 particles="dust", scrim="full")
    eb = eyebrow(slide, MX, 0.72, "The Challenge  ·  A Fair Needs A Village To Walk Down From",
                 color=C.SAFFRON_LIGHT)
    ttl = text(slide, MX, 1.16, 11.0, 0.8,
               [para("What threatens the ground", font=FE, size=30,
                     color=C.SNOW)], name="title")
    pos = grid_positions(6, 3, MX, 2.14, 3.62, 2.24, 0.24, 0.22)
    for i, ((ic, hd, bd), (x, y)) in enumerate(zip(K.CHALLENGES, pos)):
        sh = feature_card(slide, x, y, 3.62, 2.24, ic, hd, bd, head_size=15.5,
                          body_size=11, accent=C.SAFFRON, fill="14100C")
        animate_card(tl, sh, delay=800 + i * 210)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    footer(slide, _idx(), "Challenges In Preservation")
    tl.apply()
    notes(slide, K.NOTES["challenges"])


# ---- 28 future --------------------------------------------------------------
def s_future():
    slide, tl = new_slide(bg="04101F")
    scenes.scene(slide, tl, mood="dawn", rays="rays_top", mid=True,
                 ground="pines_soft", ground_h=1.6, mist="mist_low",
                 particles="pollen", scrim="full", glow=(6.667, 2.4, 7.0, "gold"))
    eb = eyebrow(slide, MX, 0.72, "The Future  ·  Money And Policy Already Exist")
    ttl = text(slide, MX, 1.16, 11.5, 0.8,
               [para("Six moves that reach the person who knows the song",
                     font=FE, size=27, color=C.SNOW)], name="title")
    pos = grid_positions(6, 3, MX, 2.14, 3.62, 2.24, 0.24, 0.22)
    for i, ((ic, hd, bd), (x, y)) in enumerate(zip(K.FUTURE, pos)):
        sh = feature_card(slide, x, y, 3.62, 2.24, ic, hd, bd, head_size=15,
                          body_size=11, accent=C.GOLD)
        animate_card(tl, sh, delay=800 + i * 210)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    footer(slide, _idx(), "Future Opportunities")
    tl.apply()
    notes(slide, K.NOTES["future"])


# ---- 29 facts ---------------------------------------------------------------
def s_facts():
    slide, tl = new_slide(bg="03110C")
    scenes.scene(slide, tl, mood="forest", rays="rays_top", mid=True,
                 ground="pines_soft", ground_h=1.4, mist="mist_low",
                 particles="fireflies", scrim="full", mandala=(10.6, 3.6, 5.6))
    eb = eyebrow(slide, MX, 0.72, "Worth Remembering  ·  Six Things About A Kauthik")
    ttl = text(slide, MX, 1.16, 11.0, 0.8,
               [para("If you take away nothing else", font=FE, size=30,
                     color=C.SNOW)], name="title")
    pos = grid_positions(6, 3, MX, 2.16, 3.62, 2.24, 0.24, 0.22)
    for i, ((ic, hd, bd), (x, y)) in enumerate(zip(K.FACTS, pos)):
        sh = feature_card(slide, x, y, 3.62, 2.24, ic, hd, bd, head_size=15.5,
                          body_size=11)
        animate_card(tl, sh, delay=750 + i * 200)
    for s in eb:
        tl.fade_in(s, delay=250)
    tl.fade_in(ttl, delay=500)
    footer(slide, _idx(), "Interesting Facts")
    tl.apply()
    notes(slide, K.NOTES["facts"])


# ---- 30 closing -------------------------------------------------------------
def s_closing():
    slide, tl = new_slide(transition="smooth", dur=D.D_TRANSITION_SLOW,
                          bg="070E1C")
    scenes.scene(slide, tl, mood="dusk", rays="rays_top", mid=True,
                 ground="temple", ground_h=2.7, mist="mist_warm",
                 particles="embers", scrim="full", flags=True,
                 glow=(6.667, 2.7, 7.5, "gold"), aipan="gold")
    eb = flanked_eyebrow(slide, 1.5, "In Closing")
    q = text(slide, 0.8, 2.5, W - 1.6, 2.2,
             [para("\u201cCulture survives when traditions are", font=FS,
                   size=42, color=C.SNOW, align="c", italic=True, line=1.16,
                   shadow=dict(blur=26, dist=4, alpha=55)),
              para("celebrated, shared, and passed on.\u201d", font=FS,
                   size=42, color=C.GOLD_LIGHT, align="c", italic=True,
                   line=1.16, shadow=dict(blur=26, dist=4, alpha=55))],
             name="closingquote")
    rl = rule(slide, W / 2 - 1.1, 4.86, 2.2, weight=1.2)
    last = text(slide, 0, 5.12, W, 0.8,
                [para("A Kauthik is not preserved by being recorded. It is "
                      "preserved by being attended.", font=FB, size=16,
                      color=C.SNOW_DIM, align="c", line=1.3)], name="lastline")
    for s in eb:
        tl.fade_in(s, delay=500, dur=1400)
    tl.rise_in(q, delay=1000, dur=2000, rise=0.03)
    tl.wipe_in(rl, delay=2400, dur=1000)
    tl.fade_in(last, delay=2700, dur=1500)
    tl.apply()
    notes(slide, K.NOTES["closing"])


# ---- 31 thanks --------------------------------------------------------------
def s_thanks():
    slide, tl = new_slide(bg="04101F")
    scenes.scene(slide, tl, mood="dawn", rays="rays_top", mid=True,
                 ground="temple", ground_h=2.4, mist="mist_low",
                 particles="dust", scrim="bottom", flags=True,
                 mandala=(6.667, 3.0, 6.2), aipan="gold")
    dev = text(slide, 0, 1.5, W, 0.7,
               [para("\u0927\u0928\u094d\u092f\u0935\u093e\u0926", font=FS,
                     size=30, color=C.GOLD_LIGHT, align="c", cs=True, track=60)],
               name="devthanks")
    big = text(slide, 0, 2.3, W, 1.5,
               [para("Thank You", font=FD, size=82, color=C.SNOW, align="c",
                     bold=True, track=40, shadow=dict(blur=30, dist=5, alpha=55))],
               name="thanks")
    rl = rule(slide, W / 2 - 1.3, 4.15, 2.6, weight=1.3)
    sub = text(slide, 0, 4.4, W, 0.6,
               [para("Questions & Discussion", font=FS, size=26,
                     color=C.GOLD_LIGHT, align="c", italic=True)], name="qd")
    tag = text(slide, 0, 5.35, W, 0.4,
               [para("Kauthik  ·  The Living Fairs Of Uttarakhand", font=FU,
                     size=12.5, color=C.SNOW_DIM, align="c", caps=True,
                     track=260)], name="tag")
    tl.fade_in(dev, delay=400, dur=1400)
    tl.rise_in(big, delay=800, dur=1700, rise=0.04)
    tl.wipe_in(rl, delay=1700, dur=900)
    tl.fade_in(sub, delay=2000, dur=1300)
    tl.fade_in(tag, delay=2400)
    tl.apply()
    notes(slide, K.NOTES["thanks"])


# ---- 32 sources -------------------------------------------------------------
def s_sources():
    slide, tl = new_slide(bg="060D1A")
    scenes.scene(slide, tl, mood="night", rays=None, mid=True, particles="dust",
                 mist="mist_low", scrim="full")
    eb = eyebrow(slide, MX, 0.7, "References & Credits")
    ttl = text(slide, MX, 1.16, 11.0, 0.6,
               [para("Sources", font=FE, size=30, color=C.SNOW)], name="title")
    half = len(K.SOURCES) // 2 + len(K.SOURCES) % 2
    colA = K.SOURCES[:half]
    colB = K.SOURCES[half:]
    sh = []
    for ci, colitems in enumerate((colA, colB)):
        x = MX + ci * 5.85
        for i, src in enumerate(colitems):
            y = 1.95 + i * 0.86
            d = dot(slide, x + 0.06, y + 0.1, 0.1, fill=C.GOLD_DEEP, name="srcdot")
            t = text(slide, x + 0.3, y, 5.35, 0.8,
                     [para(src, size=11, color=C.SNOW_DIM, line=1.24)],
                     name="src")
            sh += [d, t]
    credit = text(slide, MX, 6.55, 11.2, 0.7,
                  [para(K.CREDIT, size=10.5, color=C.SLATE_LIGHT, italic=True,
                        line=1.28, font=FS)], name="credit")
    for s in eb:
        tl.fade_in(s, delay=200)
    tl.fade_in(ttl, delay=450)
    tl.stagger(sh, start=700, gap=90, mode="fade")
    tl.fade_in(credit, delay=1800)
    footer(slide, _idx(), "Sources")
    tl.apply()
    notes(slide, K.NOTES["sources"])


# ===================================================================== driver ==
def build(force_art: bool = False):
    print("· generating background art\u2026")
    scenes.build_all(force=force_art, verbose=True)
    print("· generating maps\u2026")
    build_maps(force=force_art)

    print("· composing slides\u2026")
    s_cover()
    s_coldopen()
    s_devbhumi()
    chapter("I", "THE WORD", "What the word carries", "dusk",
            particles="dust", note_key="ch1")
    s_word()
    s_anatomy()
    chapter("II", "THE HISTORY", "A thousand years coming down the hill",
            "ember", particles="embers", ground="pines", note_key="ch2")
    s_timeline()
    s_1921()
    chapter("III", "THE LAND", "Where the fairs stand", "snow",
            particles="snowfall", note_key="ch3", rays="rays_cool")
    s_map()
    s_calendar()
    chapter("IV", "THE SOUND", "Drum, horn and the circle", "night",
            particles="fireflies", note_key="ch4", rays="rays_cool")
    s_music()
    s_dance()
    s_dance_film()
    chapter("V", "THE HANDS", "Cloth, copper, bamboo and rice paste",
            "parchment", particles="dust", note_key="ch5", rays="rays_left")
    s_attire()
    s_crafts()
    s_food()
    chapter("VI", "THE SPIRIT", "Goddess, oath and ordeal", "royal",
            particles="fireflies", note_key="ch6", rays="rays_cool")
    s_spirit()
    chapter("VII", "THE MARKETPLACE", "What the fair is worth", "parchment",
            particles="dust", ground="stalls", note_key="ch7",
            rays="rays_left")
    s_economy()
    s_tourism()
    chapter("VIII", "THE FUTURE", "Empty houses and a full ground", "dawn",
            particles="pollen", note_key="ch8")
    s_modern()
    s_challenges()
    s_future()
    s_facts()
    s_closing()
    s_thanks()
    s_sources()

    out = os.path.join(HERE, "Kauthik.pptx")
    prs.save(out)
    print(f"· saved {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides, "
          f"{os.path.getsize(out) / 1e6:.1f} MB)")
    return out


if __name__ == "__main__":
    build(force_art="--force-art" in sys.argv)
