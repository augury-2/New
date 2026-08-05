"""
KAUTHIK — PowerPoint primitives
===============================
Low-level helpers that python-pptx does not expose: character tracking, all-
caps runs, complex-script typefaces (for Devanagari), soft text shadows,
picture insertion by path with automatic part de-duplication, and shape
naming so that Morph transitions can match objects between slides.

Everything returns real, editable PowerPoint objects — there is not a single
flattened text image in the deck.
"""

from __future__ import annotations

import copy

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from design import F_BODY, F_UI, rgb

CS_FONT = "Nirmala UI"        # complex-script face for Devanagari runs


# ------------------------------------------------------------------ geometry --
def IN(v: float) -> Emu:
    return Inches(v)


# --------------------------------------------------------------------- text ---
def _rPr(run):
    rPr = run._r.get_or_add_rPr()
    return rPr


def set_tracking(run, hundredths_pt: int) -> None:
    """Character spacing. `spc` is in 1/100 pt; negative tightens."""
    _rPr(run).set("spc", str(int(hundredths_pt)))


def set_caps(run, mode: str = "all") -> None:
    _rPr(run).set("cap", mode)


def set_cs_font(run, typeface: str = CS_FONT) -> None:
    """Bind the complex-script and east-asian faces so Devanagari resolves."""
    rPr = _rPr(run)
    for tag in ("a:cs", "a:ea"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", typeface)


def set_text_shadow(run, blur: float = 26.0, dist: float = 4.0,
                    direction: int = 5400000, alpha: int = 62,
                    color: str = "000000") -> None:
    """
    A soft drop shadow on the glyphs themselves. Used sparingly on display
    type sitting over imagery — it buys real contrast without a dark box.
    """
    rPr = _rPr(run)
    old = rPr.find(qn("a:effectLst"))
    if old is not None:
        rPr.remove(old)
    xml = (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:outerShdw blurRad="%d" dist="%d" dir="%d" rotWithShape="0">'
        '<a:srgbClr val="%s"><a:alpha val="%d000"/></a:srgbClr>'
        '</a:outerShdw></a:effectLst>'
        % (int(blur * 12700), int(dist * 12700), direction, color, alpha)
    )
    from lxml import etree
    el = etree.fromstring(xml)
    # effectLst must come after ln / fill children of rPr
    rPr.append(el)


def set_line_spacing(para, mult: float) -> None:
    para.line_spacing = mult


def textbox(slide, x, y, w, h, name: str | None = None, anchor=MSO_ANCHOR.TOP,
            wrap: bool = True, autofit: bool = False):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if not autofit:
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        for tag in ("a:normAutofit", "a:spAutoFit"):
            el = bodyPr.find(qn(tag))
            if el is not None:
                bodyPr.remove(el)
    if name:
        tb._element._nvXxPr.cNvPr.set("name", name)
    return tb


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT,
         "j": PP_ALIGN.JUSTIFY}


def write(tb, paras, name: str | None = None):
    """
    Fill a textbox from a compact paragraph spec.

    Each paragraph is a dict:
      text  str | list[dict]   — a string, or runs for mixed styling
      font  str                — typeface
      size  float              — pt
      color str                — hex, no '#'
      bold / italic  bool
      align 'l'|'c'|'r'|'j'
      track int                — character spacing, 1/100 pt
      caps  bool
      line  float              — line spacing multiple
      before / after  float    — space before/after, pt
      shadow bool | dict
      cs    bool               — bind Devanagari-capable complex-script font
    """
    tf = tb.text_frame
    tf.clear()
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if "align" in spec:
            p.alignment = ALIGN[spec["align"]]
        if "line" in spec:
            p.line_spacing = spec["line"]
        if "before" in spec:
            p.space_before = Pt(spec["before"])
        if "after" in spec:
            p.space_after = Pt(spec["after"])

        runs = spec["text"] if isinstance(spec["text"], list) else [
            {k: v for k, v in spec.items() if k != "text"} | {"text": spec["text"]}]
        for rspec in runs:
            merged = {k: v for k, v in spec.items()
                      if k not in ("text", "align", "line", "before", "after")}
            merged.update(rspec)
            r = p.add_run()
            r.text = merged.get("text", "")
            f = r.font
            f.name = merged.get("font", F_BODY)
            f.size = Pt(merged.get("size", 17))
            f.bold = bool(merged.get("bold", False))
            f.italic = bool(merged.get("italic", False))
            f.color.rgb = RGBColor(*rgb(merged.get("color", "F8F6F1")))
            if merged.get("track"):
                set_tracking(r, merged["track"])
            if merged.get("caps"):
                set_caps(r)
            if merged.get("cs"):
                set_cs_font(r, merged.get("cs_font", CS_FONT))
            sh = merged.get("shadow")
            if sh:
                set_text_shadow(r, **(sh if isinstance(sh, dict) else {}))
    if name:
        tb._element._nvXxPr.cNvPr.set("name", name)
    return tb


def text(slide, x, y, w, h, paras, name=None, anchor=MSO_ANCHOR.TOP):
    return write(textbox(slide, x, y, w, h, name=name, anchor=anchor), paras)


# ------------------------------------------------------------------- shapes ---
def picture(slide, path, x, y, w=None, h=None, name=None, transparency=None):
    """
    Insert a picture. python-pptx hashes image bytes, so re-using the same
    asset across slides adds exactly one image part to the package.
    """
    pic = slide.shapes.add_picture(path, IN(x), IN(y),
                                   IN(w) if w is not None else None,
                                   IN(h) if h is not None else None)
    if name:
        pic._element._nvXxPr.cNvPr.set("name", name)
    if transparency is not None:
        set_picture_transparency(pic, transparency)
    return pic


def set_picture_transparency(pic, pct: float) -> None:
    """pct 0..1 — 0.3 means 30% transparent."""
    blip = pic._element.blipFill.find(qn("a:blip"))
    from lxml import etree
    el = etree.fromstring(
        '<a:alphaModFix xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' amt="%d"/>' % int(round((1.0 - pct) * 100000))
    )
    blip.append(el)


def rect(slide, x, y, w, h, fill=None, alpha=None, line=None, line_w=1.0,
         name=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, IN(x), IN(y), IN(w), IN(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor(*rgb(fill))
        if alpha is not None:
            _fill_alpha(sp, alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor(*rgb(line))
        sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        set_adjust(sp, radius)
    sp.text_frame.text = ""
    if name:
        sp._element._nvXxPr.cNvPr.set("name", name)
    return sp


def _fill_alpha(sp, alpha: float) -> None:
    """alpha 0..1 opacity."""
    from lxml import etree
    srgb = sp.fill.fore_color._xFill.find(qn("a:srgbClr"))
    el = etree.fromstring(
        '<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
        ' val="%d"/>' % int(round(alpha * 100000))
    )
    srgb.append(el)


def set_adjust(sp, value: float, idx: int = 0) -> None:
    """Set a preset-geometry adjustment (e.g. corner radius 0..0.5)."""
    prstGeom = sp._element.spPr.find(qn("a:prstGeom"))
    avLst = prstGeom.find(qn("a:avLst"))
    if avLst is None:
        from lxml import etree
        avLst = etree.SubElement(prstGeom, qn("a:avLst"))
    for gd in list(avLst):
        avLst.remove(gd)
    from lxml import etree
    gd = etree.SubElement(avLst, qn("a:gd"))
    gd.set("name", "adj%d" % (idx + 1) if idx else "adj")
    gd.set("fmla", "val %d" % int(round(value * 100000)))


def hairline(slide, x, y, w, color="E3B34A", weight=0.9, alpha=None, name=None):
    """A thin horizontal rule — the deck's main structural device."""
    return rect(slide, x, y, w, weight / 72.0, fill=color, alpha=alpha, name=name)


def vline(slide, x, y, h, color="E3B34A", weight=0.9, alpha=None, name=None):
    return rect(slide, x, y, weight / 72.0, h, fill=color, alpha=alpha, name=name)


def dot(slide, cx, cy, d, fill="E3B34A", alpha=None, line=None, line_w=1.0,
        name=None):
    return rect(slide, cx - d / 2, cy - d / 2, d, d, fill=fill, alpha=alpha,
                line=line, line_w=line_w, name=name, shape=MSO_SHAPE.OVAL)


def ring(slide, cx, cy, d, color="E3B34A", weight=1.0, name=None):
    return rect(slide, cx - d / 2, cy - d / 2, d, d, fill=None, line=color,
                line_w=weight, name=name, shape=MSO_SHAPE.OVAL)


# ------------------------------------------------------------------- notes ----
def notes(slide, body: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.text = body.strip()
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11)
            r.font.name = F_UI


# --------------------------------------------------------------- z-ordering ---
def send_to_back(shape) -> None:
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    # insert after the two required leading elements (nvGrpSpPr, grpSpPr)
    spTree.insert(2, shape._element)


def bring_to_front(shape) -> None:
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.append(shape._element)


# ------------------------------------------------------------ slide plumbing ---
def blank_slide(prs):
    layout = prs.slide_layouts[6]          # Blank
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def set_slide_bg(slide, color: str) -> None:
    from lxml import etree
    cSld = slide._element.find(qn("p:cSld"))
    old = cSld.find(qn("p:bg"))
    if old is not None:
        cSld.remove(old)
    xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="%s"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>' % color
    )
    cSld.insert(0, etree.fromstring(xml))
